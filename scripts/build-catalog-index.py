#!/usr/bin/env python3
"""
Build Architecture Catalog Index

Reads system_architecture_catalog.pptx and outputs architecture_catalog.json
with metadata for each slide: index, title, subtitle, kind (section vs diagram),
tags, talk track preview, and section grouping.

Usage:
    python build-catalog-index.py
"""

import json
import re
import sys
import unicodedata
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("Error: python-pptx is required. Install with: pip3 install python-pptx")
    sys.exit(1)

try:
    from lxml import etree
except ImportError:
    print("Error: lxml is required. Install with: pip3 install lxml")
    sys.exit(1)

SCRIPT_DIR = Path(__file__).parent
SKILL_DIR = SCRIPT_DIR.parent
CATALOG_PPTX = SKILL_DIR / "assets" / "databricks" / "system_architecture_catalog.pptx"
OUTPUT_JSON = SKILL_DIR / "assets" / "databricks" / "architecture_catalog.json"

# Slides to skip (1-based): slide 1 = disclaimer, slide 2 = changelog
SKIP_SLIDES = {1, 2}

# Section header: low shape count + typically a title-only layout
SECTION_MAX_SHAPES = 4


def extract_tags_and_talk_track(slide):
    """Extract tags and talk track from slide notes.

    Tags are lines starting with 'tag:' (comma-separated values).
    Talk track is everything else in the notes.
    """
    tags = []
    talk_track_lines = []

    try:
        if not slide.has_notes_slide:
            return tags, ""
        notes_text = slide.notes_slide.notes_text_frame.text
    except Exception:
        return tags, ""

    if not notes_text:
        return tags, ""

    for line in notes_text.split("\n"):
        stripped = line.strip()
        # Match "tag:" or "tag: " at start of line
        if re.match(r'^tags?\s*:', stripped, re.IGNORECASE):
            tag_part = re.sub(r'^tags?\s*:\s*', '', stripped, flags=re.IGNORECASE)
            for t in tag_part.split(","):
                t = t.strip()
                if t:
                    tags.append(t)
        elif stripped:
            talk_track_lines.append(stripped)

    # Also check for standalone level tags like "L200" at start of notes
    if not tags and talk_track_lines:
        first = talk_track_lines[0]
        if re.match(r'^L\d{3}', first):
            tags.append(first)
            talk_track_lines = talk_track_lines[1:]

    talk_track = " ".join(talk_track_lines)
    return tags, talk_track


def get_slide_title(slide):
    """Extract the title text from a slide."""
    if slide.shapes.title:
        return slide.shapes.title.text.strip()
    # Fallback: look for the first text shape
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return shape.text_frame.text.strip()[:100]
    return ""


def get_slide_subtitle(slide):
    """Extract subtitle text (second text element) from a slide."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            texts.append(shape.text_frame.text.strip())
    if len(texts) >= 2:
        # First is title, second might be subtitle
        candidate = texts[1]
        # Only return if it's short enough to be a subtitle
        if len(candidate) < 200:
            return candidate
    return ""


def classify_slide(slide):
    """Classify slide as 'section' or 'diagram'."""
    shape_count = len(slide.shapes)
    if shape_count <= SECTION_MAX_SHAPES:
        return "section"
    return "diagram"


# EMU to inches conversion factor
EMU_PER_INCH = 914400

# XML namespaces used in OOXML
NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'


def _normalize_text(text):
    """Normalize Unicode whitespace and control characters for comparison."""
    # Replace non-breaking spaces, vertical tabs, etc. with regular spaces
    text = unicodedata.normalize('NFKC', text)
    text = re.sub(r'[\u00a0\u200b\u200c\u200d\ufeff]', ' ', text)
    text = re.sub(r'\s+', ' ', text)
    return text.strip()


def _emu_to_inches(emu):
    """Convert EMU (English Metric Units) to inches, rounded to 2 decimal places."""
    if emu is None:
        return 0.0
    return round(int(emu) / EMU_PER_INCH, 2)


def _get_group_transform(grpSp_elem):
    """Extract group coordinate transform from a grpSp element.

    Returns a dict with off_x/y, ext_cx/cy, chOff_x/y, chExt_cx/cy,
    has_rotation, or None if transform not found.
    """
    for ns in [NS_P, NS_A]:
        grpSpPr = grpSp_elem.find(f'{{{ns}}}grpSpPr')
        if grpSpPr is not None:
            xfrm = grpSpPr.find(f'{{{NS_A}}}xfrm')
            if xfrm is not None:
                off = xfrm.find(f'{{{NS_A}}}off')
                ext = xfrm.find(f'{{{NS_A}}}ext')
                chOff = xfrm.find(f'{{{NS_A}}}chOff')
                chExt = xfrm.find(f'{{{NS_A}}}chExt')
                if all(e is not None for e in [off, ext, chOff, chExt]):
                    return {
                        "off_x": int(off.get('x', '0')),
                        "off_y": int(off.get('y', '0')),
                        "ext_cx": int(ext.get('cx', '0')),
                        "ext_cy": int(ext.get('cy', '0')),
                        "chOff_x": int(chOff.get('x', '0')),
                        "chOff_y": int(chOff.get('y', '0')),
                        "chExt_cx": int(chExt.get('cx', '0')),
                        "chExt_cy": int(chExt.get('cy', '0')),
                        "has_rotation": xfrm.get('rot') is not None,
                    }
    return None


def _child_to_slide_coords(child_x, child_y, gxfrm):
    """Convert child-space EMU coordinates to slide-space EMU.

    Uses: slide = off + (child - chOff) * (ext / chExt)
    """
    if gxfrm["chExt_cx"] == 0 or gxfrm["chExt_cy"] == 0:
        return child_x, child_y

    slide_x = gxfrm["off_x"] + int(
        (child_x - gxfrm["chOff_x"]) * gxfrm["ext_cx"] / gxfrm["chExt_cx"]
    )
    slide_y = gxfrm["off_y"] + int(
        (child_y - gxfrm["chOff_y"]) * gxfrm["ext_cy"] / gxfrm["chExt_cy"]
    )
    return slide_x, slide_y


def _extract_shape_text(sp_elem):
    """Extract concatenated text from a shape element's txBody.

    txBody may be in p: or a: namespace depending on context.
    Paragraphs and runs inside are always in a: namespace.
    """
    txBody = sp_elem.find(f'{{{NS_P}}}txBody')
    if txBody is None:
        txBody = sp_elem.find(f'{{{NS_A}}}txBody')
    if txBody is None:
        return ""
    texts = []
    for p in txBody.findall(f'{{{NS_A}}}p'):
        runs = p.findall(f'{{{NS_A}}}r')
        para_text = "".join(
            r.findtext(f'{{{NS_A}}}t', default='') for r in runs
        )
        if para_text.strip():
            texts.append(para_text.strip())
    return " ".join(texts)


def _get_shape_font_size(sp_elem):
    """Get the font size (in pt) from the first run of a shape, or None.

    Searches for rPr inside a:r (run properties) or a:pPr/a:defRPr
    (default paragraph run properties). txBody may be in p: or a: namespace.
    """
    # Find txBody in either namespace
    txBody = sp_elem.find(f'{{{NS_P}}}txBody')
    if txBody is None:
        txBody = sp_elem.find(f'{{{NS_A}}}txBody')
    if txBody is None:
        return None

    rPr = txBody.find(f'.//{{{NS_A}}}r/{{{NS_A}}}rPr')
    if rPr is not None:
        sz = rPr.get('sz')
        if sz:
            return round(int(sz) / 100)  # hundredths of a point -> pt
    defRPr = txBody.find(f'.//{{{NS_A}}}pPr/{{{NS_A}}}defRPr')
    if defRPr is not None:
        sz = defRPr.get('sz')
        if sz:
            return round(int(sz) / 100)
    return None


def _get_shape_position(sp_elem, group_xfrm=None):
    """Extract position and size from a shape's spPr/xfrm or grpSpPr/xfrm.

    spPr may be in p: namespace (top-level shapes) with a:xfrm inside,
    or in a: namespace (inside groups). Try both.

    If group_xfrm is provided, converts child-space coordinates to slide-space
    and scales width/height by the group's scaling factor.
    """
    for ns in [NS_P, NS_A]:
        for xfrm_parent_tag in ['spPr', 'grpSpPr']:
            parent = sp_elem.find(f'{{{ns}}}{xfrm_parent_tag}')
            if parent is not None:
                xfrm = parent.find(f'{{{NS_A}}}xfrm')
                if xfrm is not None:
                    off = xfrm.find(f'{{{NS_A}}}off')
                    ext = xfrm.find(f'{{{NS_A}}}ext')
                    if off is not None and ext is not None:
                        child_x = int(off.get('x', '0'))
                        child_y = int(off.get('y', '0'))
                        child_cx = int(ext.get('cx', '0'))
                        child_cy = int(ext.get('cy', '0'))

                        if group_xfrm is not None:
                            slide_x, slide_y = _child_to_slide_coords(
                                child_x, child_y, group_xfrm
                            )
                            # Scale dimensions by group scaling factor
                            sx = (group_xfrm["ext_cx"] / group_xfrm["chExt_cx"]
                                  if group_xfrm["chExt_cx"] else 1.0)
                            sy = (group_xfrm["ext_cy"] / group_xfrm["chExt_cy"]
                                  if group_xfrm["chExt_cy"] else 1.0)
                            return {
                                "left": _emu_to_inches(slide_x),
                                "top": _emu_to_inches(slide_y),
                                "width": _emu_to_inches(int(child_cx * sx)),
                                "height": _emu_to_inches(int(child_cy * sy)),
                            }
                        else:
                            return {
                                "left": _emu_to_inches(child_x),
                                "top": _emu_to_inches(child_y),
                                "width": _emu_to_inches(child_cx),
                                "height": _emu_to_inches(child_cy),
                            }
    return None


def _get_shape_name(sp_elem):
    """Extract shape name from nvSpPr/cNvPr."""
    cNvPr = sp_elem.find(f'{{{NS_P}}}nvSpPr/{{{NS_P}}}cNvPr')
    if cNvPr is None:
        # Try drawingml namespace (inside groups)
        cNvPr = sp_elem.find(f'{{{NS_P}}}nvSpPr/{{{NS_A}}}cNvPr')
    if cNvPr is None:
        # Generic fallback: look for any cNvPr descendant
        cNvPr = sp_elem.find(f'.//{{{NS_P}}}cNvPr')
    if cNvPr is not None:
        return cNvPr.get('name', '')
    return ''


def _should_skip_text(text):
    """Return True if this text should be excluded from the label index."""
    if not text:
        return True
    if len(text) <= 1:
        return True
    # Slide number placeholder
    if text.strip() in ('‹#›', '<#>', '<<#>>'):
        return True
    return False


def extract_text_labels(slide):
    """Extract all text-bearing shapes from a slide's XML, including groups.

    Returns a list of dicts with: text, shape_name, in_group, position, font_size.
    """
    cSld = slide._element.find(f'{{{NS_P}}}cSld')
    if cSld is None:
        return []

    labels = []

    def _process_shapes(parent, in_group=False, group_xfrm=None):
        # Process top-level shapes (p:sp)
        for sp in parent.findall(f'{{{NS_P}}}sp'):
            text = _extract_shape_text(sp)
            text = _normalize_text(text)
            if _should_skip_text(text):
                continue

            label = {
                "text": text[:100],
                "shape_name": _get_shape_name(sp),
                "in_group": in_group,
            }
            pos = _get_shape_position(sp, group_xfrm=group_xfrm)
            if pos:
                label["position"] = pos
            font_size = _get_shape_font_size(sp)
            if font_size:
                label["font_size"] = font_size

            labels.append(label)

        # Recurse into groups (p:grpSp)
        for grpSp in parent.findall(f'{{{NS_P}}}grpSp'):
            gxfrm = _get_group_transform(grpSp)
            _process_shapes(grpSp, in_group=True, group_xfrm=gxfrm)

    spTree = cSld.find(f'{{{NS_P}}}spTree')
    if spTree is not None:
        _process_shapes(spTree)

    return labels


# Reference Architecture slides (65-76) share a common base diagram.
# These are the only slides targeted for customer-specific modification.
REFERENCE_ARCH_RANGE = range(65, 77)

# Known Databricks product names — these rarely need replacement.
_DATABRICKS_PRODUCTS = {
    'Unity Catalog', 'Delta Lake', 'Iceberg', 'Spark', 'Photon',
    'MLFlow', 'Mosaic AI', 'Lakeflow', 'Lakeflow Connect',
    'Lakeflow Jobs', 'Lakebase', 'AutoLoader', 'Delta Sharing',
    'Databricks SQL', 'Databricks Apps', 'Databricks', 'Genie',
    'Dashboards', 'Marketplace', 'Clean Rooms', 'Asset Bundles',
    'Terraform Provider', 'SDKs', 'AI Gateway', 'Vector Search',
    'Feature Serving', 'Feature Engineering', 'Model Serving',
    'Agent Framework', 'Agent Bricks', 'Predictive IO',
    'Predictive optimization', 'Assistant', 'AI Functions',
    'Spark Declarative Pipelines', 'Spark /Photon',
    'Connectors and APIs', 'Traditional ML',
}

# Flow titles for each reference architecture slide (slide_num -> highlight description)
_REFERENCE_FLOW_TITLES = {
    65: "Full platform overview — no specific flow highlighted",
    66: "Built-in ingestion from SaaS and databases (Lakeflow Connect)",
    67: "Batch ingestion and ETL (AutoLoader, Spark, Declarative Pipelines)",
    68: "Streaming and Change Data Capture",
    69: "Machine Learning — traditional ML workflow",
    70: "Generative AI — Agents and LLM serving",
    71: "Business Intelligence (SQL Warehouses, Dashboards, Genie)",
    72: "Business Apps (Databricks Apps)",
    73: "Lakehouse Federation (query external data in place)",
    74: "Catalog Federation (integrate external HMS/catalogs)",
    75: "Sharing Data outbound (Delta Sharing)",
    76: "Consuming Shared Data inbound",
}


def build_reference_arch_guide(text_labels, slide_num):
    """Build a customization guide for a reference architecture slide.

    Categorizes labels into: sources, consumers, and replaceable labels.
    Adds a prose customization_guide explaining what to modify.
    """
    sources = []
    consumers = []
    replaceable = []

    for label in text_labels:
        text = label['text']
        pos = label.get('position', {})
        left = pos.get('left', 5)
        in_group = label.get('in_group', False)

        # Data sources: grouped shapes on the left edge
        if in_group and left < 2.5:
            # Skip tiny utility labels
            if len(text) > 3:
                sources.append(text)
            continue

        # External consumers: grouped shapes on the right edge
        if in_group and left > 9.5:
            if len(text) > 3:
                consumers.append(text)
            continue

        # Skip known Databricks products
        # Normalize for matching: strip parens, compare core text
        core = text.split('(')[0].strip()
        if core in _DATABRICKS_PRODUCTS:
            continue

        # Skip very short labels and structural labels
        if len(text) <= 3:
            continue

    # Deduplicate
    sources = list(dict.fromkeys(sources))
    consumers = list(dict.fromkeys(consumers))

    flow_desc = _REFERENCE_FLOW_TITLES.get(slide_num, "")

    guide = {
        "flow": flow_desc,
        "sources": sources,
        "consumers": consumers,
        "customization_guide": (
            "Replace 'sources' with the customer's actual data systems "
            "(e.g., 'Epic EHR', 'Kafka', 'SAP'). Replace 'consumers' with "
            "their downstream tools (e.g., 'Tableau', 'Internal App'). "
            "Databricks product names inside the platform should usually "
            "stay as-is. Use the title override to brand the slide."
        ),
    }

    return guide


def build_catalog():
    if not CATALOG_PPTX.exists():
        print(f"Error: Catalog not found at {CATALOG_PPTX}")
        sys.exit(1)

    prs = Presentation(str(CATALOG_PPTX))
    total_slides = len(prs.slides)
    print(f"Reading {total_slides} slides from {CATALOG_PPTX.name}")

    slides = []
    current_section = None
    section_order = 0

    for idx, slide in enumerate(prs.slides):
        slide_num = idx + 1  # 1-based

        if slide_num in SKIP_SLIDES:
            continue

        title = get_slide_title(slide)
        subtitle = get_slide_subtitle(slide)
        kind = classify_slide(slide)
        tags, talk_track = extract_tags_and_talk_track(slide)
        shape_count = len(slide.shapes)

        # Track sections
        if kind == "section":
            current_section = title
            section_order += 1

        entry = {
            "slide": slide_num,
            "title": title,
            "kind": kind,
            "tags": tags,
            "shape_count": shape_count,
        }

        if subtitle and subtitle != title:
            entry["subtitle"] = subtitle

        if talk_track:
            # Truncate talk track preview to 300 chars
            entry["talk_track"] = talk_track[:300] + ("..." if len(talk_track) > 300 else "")

        if current_section and kind != "section":
            entry["section"] = current_section

        # Extract text labels for diagram slides
        if kind == "diagram":
            text_labels = extract_text_labels(slide)
            if text_labels:
                entry["text_labels"] = text_labels

            # Add customization guide for reference architecture slides
            if slide_num in REFERENCE_ARCH_RANGE and text_labels:
                entry["reference_arch"] = build_reference_arch_guide(text_labels, slide_num)

        slides.append(entry)

    # Build section index
    sections = []
    seen_sections = {}
    for s in slides:
        if s["kind"] == "section":
            if s["title"] not in seen_sections:
                seen_sections[s["title"]] = len(sections)
                sections.append({
                    "title": s["title"],
                    "slide": s["slide"],
                    "diagrams": []
                })
        elif s["kind"] == "diagram" and s.get("section"):
            sec_idx = seen_sections.get(s["section"])
            if sec_idx is not None:
                sections[sec_idx]["diagrams"].append(s["slide"])

    catalog = {
        "source": CATALOG_PPTX.name,
        "total_slides": total_slides,
        "importable_slides": len([s for s in slides if s["kind"] == "diagram"]),
        "sections": sections,
        "slides": slides,
    }

    OUTPUT_JSON.parent.mkdir(parents=True, exist_ok=True)
    with open(OUTPUT_JSON, "w", encoding="utf-8") as f:
        json.dump(catalog, f, indent=2, ensure_ascii=False)

    diagram_count = catalog["importable_slides"]
    section_count = len(sections)
    print(f"Wrote {OUTPUT_JSON}")
    print(f"  Sections: {section_count}")
    print(f"  Diagrams: {diagram_count}")
    print(f"  Total indexed: {len(slides)} (skipped {len(SKIP_SLIDES)} meta slides)")


if __name__ == "__main__":
    build_catalog()
