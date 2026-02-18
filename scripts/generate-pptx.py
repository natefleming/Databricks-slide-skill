#!/usr/bin/env python3
"""
Databricks Slide Deck Generator (Template-Based)

Generates PowerPoint presentations using the official Databricks corporate template.
Output can be imported directly into Google Slides.

Usage:
    python generate-pptx.py --input content.json --output presentation.pptx
"""

import json
import argparse
import re
import sys
import unicodedata
from copy import deepcopy
from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError as e:
    print(f"Error: python-pptx is required. Install with: pip3 install python-pptx")
    print(f"Details: {e}")
    sys.exit(1)

# =============================================================================
# Constants
# =============================================================================

SCRIPT_DIR = Path(__file__).parent
SKILL_DIR = SCRIPT_DIR.parent
TEMPLATE_PATH = SKILL_DIR / "assets" / "databricks" / "template.pptx"
THEME_PATH = SKILL_DIR / "themes" / "databricks.json"
CATALOG_PPTX_PATH = SKILL_DIR / "assets" / "databricks" / "system_architecture_catalog.pptx"
CATALOG_JSON_PATH = SKILL_DIR / "assets" / "databricks" / "architecture_catalog.json"
ICON_CATALOG_PATH = SKILL_DIR / "assets" / "icons" / "icon_catalog.json"
ICON_DIR = SKILL_DIR / "assets" / "icons"

# Valid slide types (26 total)
VALID_SLIDE_TYPES = {
    # Existing types (17)
    "title", "section", "content", "two-column", "three-column",
    "big-number", "callout", "quote", "closing",
    "agenda", "timeline", "icon-grid", "stat-row", "pros-cons",
    "comparison", "checklist", "logos",
    # New types (8)
    "two-column-icons", "three-column-icons", "cards",
    "card-right", "card-left", "card-full",
    "one-column", "section-description",
    # Imported slides
    "architecture",
}

# Layout name mappings (our type -> template layout name patterns)
# For slides with prefer_dark=True, get_layout searches dark_layouts first
LAYOUT_MAPPINGS = {
    # Structural slides (dark by default) - use Databricks dark templates
    "title": ["1_3 Title Slide B - Dark", "3 Title Slide B - Light", "TITLE"],
    "section": ["Content E - Power Statement 3", "SECTION_HEADER"],  # Dark statement layout
    "callout": ["Content E - Power Statement 2_1", "MAIN_POINT"],  # Dark with title+subtitle
    "quote": ["Content E - Power Statement 2_1", "MAIN_POINT"],  # Dark with title+subtitle
    "closing": ["Z - Closing Dark", "Z - Closing Light"],
    # Content slides (light by default)
    "content": ["7 Content A - Basic", "TITLE_AND_BODY"],
    "two-column": ["9 Content B - 2 Column", "TITLE_AND_TWO_COLUMNS"],
    "three-column": ["11 Content C - 3 Column"],
    "big-number": ["Content E - Power Statement 1", "BIG_NUMBER"],
    # New template types (these are already on Master 1/2)
    "two-column-icons": ["10 Content B - 2 Column w/ Icon Spot"],
    "three-column-icons": ["12 Content C - 3 Column w/ Icon Spot"],
    "cards": ["13 Content C - 3 Column Cards"],
    "card-right": ["14 Content D - Card Right"],
    "card-left": ["15 Content D - Card Left"],
    "card-full": ["16 Content D - Card Large"],
    "one-column": ["7 Content A - Basic", "ONE_COLUMN_TEXT"],  # Master 1/2 for footer
    "section-description": ["Content E - Power Statement 2", "SECTION_TITLE_AND_DESCRIPTION"],  # Master 1 for footer
    # Hybrid types (use CUSTOM for clean slate with footer from master)
    "agenda": ["CUSTOM"],
    "timeline": ["CUSTOM"],
    "icon-grid": ["CUSTOM"],
    "stat-row": ["CUSTOM"],
    "pros-cons": ["CUSTOM"],
    "comparison": ["CUSTOM"],
    "checklist": ["CUSTOM"],
    "logos": ["CUSTOM"],
}


def load_colors_from_theme(theme_path: Path = THEME_PATH) -> Dict[str, str]:
    """Load color palette from theme JSON file.

    Falls back to hardcoded Databricks brand colors if theme file
    is missing or invalid.
    """
    try:
        with open(theme_path, 'r') as f:
            theme = json.load(f)

        return {
            "accent": theme["modes"]["light"]["accent"],
            "dark_bg": theme["modes"]["dark"]["background"],
            "light_bg": theme["modes"]["light"]["background"],
            "text_dark": theme["modes"]["light"]["text_primary"],
            "text_light": theme["modes"]["dark"]["text_primary"],
            "text_secondary": theme["modes"]["light"]["text_secondary"],
            "green": theme["elements"]["pros_header_color"],
            "red": theme["elements"]["cons_header_color"],
            "divider": theme["elements"]["stat_row_divider"],
        }
    except FileNotFoundError:
        print(f"Warning: Theme file not found at {theme_path}, using defaults")
    except json.JSONDecodeError as e:
        print(f"Warning: Invalid JSON in theme file: {e}, using defaults")
    except KeyError as e:
        print(f"Warning: Missing key in theme file: {e}, using defaults")

    # Fallback to hardcoded Databricks brand colors
    return {
        "accent": "#FF3621",
        "dark_bg": "#1B3139",
        "light_bg": "#F5F3F0",
        "text_dark": "#1B3139",
        "text_light": "#FFFFFF",
        "text_secondary": "#6B7280",
        "green": "#10B981",
        "red": "#EF4444",
        "divider": "#E5E7EB",
    }


# Colors loaded from theme file (with hardcoded fallback)
COLORS = load_colors_from_theme()


def load_font_from_theme(theme_path: Path = THEME_PATH) -> str:
    """Load font family from theme JSON file."""
    try:
        with open(theme_path, 'r') as f:
            theme = json.load(f)
        return theme.get("typography", {}).get("font_family", "DM Sans")
    except (FileNotFoundError, json.JSONDecodeError, KeyError):
        return "DM Sans"


# Font family from theme (with fallback)
FONT_FAMILY = load_font_from_theme()

_color_cache: Dict[str, RGBColor] = {}


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor with caching."""
    if hex_color in _color_cache:
        return _color_cache[hex_color]
    clean_hex = hex_color.lstrip('#')
    color = RGBColor(int(clean_hex[0:2], 16), int(clean_hex[2:4], 16), int(clean_hex[4:6], 16))
    _color_cache[hex_color] = color
    return color


# Precompiled regex for accent text parsing
_ACCENT_PATTERN = re.compile(r'\*([^*]+)\*')


def parse_accent_text(text: str) -> List[Tuple[str, bool]]:
    """Parse text for *accent* markers.

    Returns list of (text, is_accent) tuples.
    Example: "Hello *world* today" -> [("Hello ", False), ("world", True), (" today", False)]
    """
    if '*' not in text:
        return [(text, False)]

    segments = []
    last_end = 0

    for match in _ACCENT_PATTERN.finditer(text):
        # Add text before this match (if any)
        if match.start() > last_end:
            segments.append((text[last_end:match.start()], False))
        # Add the accented text (without asterisks)
        segments.append((match.group(1), True))
        last_end = match.end()

    # Add remaining text after last match
    if last_end < len(text):
        segments.append((text[last_end:], False))

    return segments if segments else [(text, False)]


# =============================================================================
# Generator Class
# =============================================================================

class DatabricksSlideGenerator:
    """Generate Databricks-branded PowerPoint presentations using corporate template."""

    def __init__(self, template_path: Path = TEMPLATE_PATH):
        if not template_path.exists():
            raise FileNotFoundError(f"Template not found: {template_path}")

        # Load template
        self.prs = Presentation(str(template_path))
        self.slide_count = 0

        # Layout caches for background mode support:
        #
        # - layouts: default cache, prefers light (first occurrence wins during iteration)
        # - light_layouts: explicit light-background layouts only (from light masters)
        # - dark_layouts: explicit dark-background layouts only (from dark masters)
        #
        # Current usage:
        #   - Template slides use `layouts` (defaults to light)
        #   - Hybrid slides use `dark_layouts` via prefer_dark=True
        #
        # Future: light_layouts enables explicit "prefer_light=True" for user-selected
        # light/dark mode per slide. Keeping all three caches for this extensibility.
        self.layouts: Dict[str, Any] = {}
        self.light_layouts: Dict[str, Any] = {}
        self.dark_layouts: Dict[str, Any] = {}

        for master in self.prs.slide_masters:
            is_dark_master = self._is_dark_background(master)

            for layout in master.slide_layouts:
                # Store in light/dark specific cache
                if is_dark_master:
                    self.dark_layouts[layout.name] = layout
                else:
                    self.light_layouts[layout.name] = layout

                # Store in main cache only if not already present (prefer first/light)
                if layout.name not in self.layouts:
                    self.layouts[layout.name] = layout

        # Clear existing slides (keep layouts)
        self._clear_slides()

    def _clear_slides(self) -> None:
        """Remove all slides from the presentation while keeping layouts."""
        # Delete slides from end to start to avoid index issues
        for i in range(len(self.prs.slides) - 1, -1, -1):
            rId = self.prs.slides._sldIdLst[i].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[i]

    def _is_dark_background(self, master) -> bool:
        """Check if a slide master has a dark background color."""
        try:
            fill = master.background.fill
            if fill.type is not None:
                if hasattr(fill, 'fore_color') and hasattr(fill.fore_color, 'rgb'):
                    rgb = fill.fore_color.rgb
                    if rgb:
                        rgb_str = str(rgb).upper()
                        # Known dark color (Databricks dark bg)
                        if rgb_str == '1B3139':
                            return True
                        # General darkness check: RGB sum < 384 means avg < 128
                        try:
                            r = int(rgb_str[0:2], 16)
                            g = int(rgb_str[2:4], 16)
                            b = int(rgb_str[4:6], 16)
                            if (r + g + b) < 384:
                                return True
                        except (ValueError, IndexError):
                            pass
        except (AttributeError, TypeError):
            pass
        return False

    def get_layout(self, slide_type: str, prefer_dark: bool = False):
        """Get the best matching layout for a slide type.

        Args:
            slide_type: The type of slide to get a layout for
            prefer_dark: If True, prefer dark background layouts when available
        """
        patterns = LAYOUT_MAPPINGS.get(slide_type, ["BLANK"])

        # Choose which caches to search based on preference
        if prefer_dark:
            search_caches = [self.dark_layouts, self.layouts]
        else:
            search_caches = [self.layouts]

        for pattern in patterns:
            for cache in search_caches:
                # Try exact match first
                if pattern in cache:
                    return cache[pattern]
            for cache in search_caches:
                # Try partial match
                for name, layout in cache.items():
                    if pattern in name:
                        return layout

        # Fallback to BLANK
        return self.layouts.get("BLANK", list(self.layouts.values())[0])

    def get_placeholder(self, slide, idx: int = None, ph_type: int = None):
        """Get placeholder by index or type."""
        for shape in slide.placeholders:
            if idx is not None and shape.placeholder_format.idx == idx:
                return shape
            if ph_type is not None and shape.placeholder_format.type == ph_type:
                return shape
        return None

    def get_placeholders_by_type(self, slide, ph_type: int) -> List:
        """Get all placeholders of a certain type, sorted by position (left to right, top to bottom)."""
        matching = []
        for shape in slide.placeholders:
            if shape.placeholder_format.type == ph_type:
                matching.append(shape)
        # Sort by top position first, then left
        matching.sort(key=lambda s: (s.top, s.left))
        return matching

    def fill_text(self, placeholder, text: str, font_size: int = None,
                  bold: bool = None, color: str = None) -> None:
        """Fill a placeholder with styled text.

        Supports accent text: wrap words in *asterisks* to highlight
        them in the accent color (Databricks orange).
        """
        if placeholder is None:
            return

        tf = placeholder.text_frame
        tf.clear()
        p = tf.paragraphs[0]

        # Parse for accent markers
        segments = parse_accent_text(text)
        has_accent = any(is_accent for _, is_accent in segments)

        if has_accent:
            # Use runs for mixed formatting
            for i, (segment_text, is_accent) in enumerate(segments):
                if i == 0:
                    run = p.runs[0] if p.runs else p.add_run()
                else:
                    run = p.add_run()
                run.text = segment_text

                if font_size:
                    run.font.size = Pt(font_size)
                if bold is not None:
                    run.font.bold = bold

                # Apply accent color or specified color
                if is_accent:
                    run.font.color.rgb = hex_to_rgb(COLORS["accent"])
                elif color:
                    run.font.color.rgb = hex_to_rgb(color)
        else:
            # Simple case: no accent markers
            p.text = text
            if font_size:
                p.font.size = Pt(font_size)
            if bold is not None:
                p.font.bold = bold
            if color:
                p.font.color.rgb = hex_to_rgb(color)

    def fill_bullets(self, placeholder, items: List[str], font_size: int = None) -> None:
        """Fill a placeholder with bullet points."""
        if placeholder is None or not items:
            return

        tf = placeholder.text_frame
        tf.clear()

        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.level = 0
            if font_size:
                p.font.size = Pt(font_size)

    def add_textbox(self, slide, text: str, left: float, top: float,
                    width: float, height: float, font_size: int = 18,
                    bold: bool = False, color: str = None,
                    alignment: int = None) -> None:
        """Add a styled textbox to a slide."""
        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = FONT_FAMILY
        p.font.size = Pt(font_size)
        p.font.bold = bold
        if color:
            p.font.color.rgb = hex_to_rgb(color)
        if alignment:
            p.alignment = alignment

    def _create_slide(self, slide_type: str, data: Dict[str, Any], prefer_dark: bool = False):
        """Create slide with layout and handle common setup.

        Handles:
        - Incrementing slide count
        - Getting the appropriate layout (light or dark)
        - Adding the slide to the presentation
        - Adding speaker notes if present in data

        Args:
            slide_type: The type of slide (maps to LAYOUT_MAPPINGS)
            data: Slide data dict, may contain "notes" key
            prefer_dark: If True, prefer dark-background layout

        Returns:
            The created slide object
        """
        self.slide_count += 1
        layout = self.get_layout(slide_type, prefer_dark=prefer_dark)
        slide = self.prs.slides.add_slide(layout)

        # Handle speaker notes for all slide types
        if data.get("notes"):
            slide.notes_slide.notes_text_frame.text = data["notes"]

        return slide

    # =========================================================================
    # Direct Template Slides
    # =========================================================================

    def add_title_slide(self, data: Dict[str, Any], prefer_dark: bool = True) -> None:
        """Create title slide using template layout.

        Args:
            data: Slide data dict
            prefer_dark: Use dark background (default True for structural slides)
        """
        from lxml import etree
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        slide = self._create_slide("title", data, prefer_dark=prefer_dark)

        # Fill title (idx 0)
        # The layout's title placeholder is 56pt with noAutofit and anchor=b.
        # Long or multi-line titles at 56pt overflow upward into the
        # Databricks logo. Use a smaller font that fits the box, and
        # explicitly set anchor=b + normAutofit at the slide level.
        title_ph = self.get_placeholder(slide, idx=0)
        title_text = data.get("title", "Presentation Title")
        title_font_size = 44 if "\n" in title_text or len(title_text) > 30 else 56
        self.fill_text(title_ph, title_text, font_size=title_font_size)

        if title_ph is not None:
            ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
            txBody = title_ph._element.find(f'{{{ns_p}}}txBody')
            if txBody is None:
                txBody = title_ph._element.find(f'{{{ns_a}}}txBody')
            if txBody is not None:
                bodyPr = txBody.find(f'{{{ns_a}}}bodyPr')
                if bodyPr is not None:
                    bodyPr.set('anchor', 'b')
                    for child in list(bodyPr):
                        localname = etree.QName(child.tag).localname
                        if localname in ('noAutofit', 'spAutoFit', 'normAutofit'):
                            bodyPr.remove(child)
                    etree.SubElement(bodyPr, f'{{{ns_a}}}normAutofit')

        # Fill subtitle (idx 1) - may contain author/date
        subtitle_ph = self.get_placeholder(slide, idx=1)
        subtitle_parts = []
        if data.get("subtitle"):
            subtitle_parts.append(data["subtitle"])
        self.fill_text(subtitle_ph, "\n".join(subtitle_parts) if subtitle_parts else "")

        # Author/date in second subtitle if available (idx 2)
        author_ph = self.get_placeholder(slide, idx=2)
        if author_ph:
            author_parts = []
            if data.get("author"):
                author_parts.append(data["author"])
            if data.get("date"):
                author_parts.append(data["date"])
            self.fill_text(author_ph, " | ".join(author_parts) if author_parts else "")

    def add_section_slide(self, data: Dict[str, Any], prefer_dark: bool = True) -> None:
        """Create section divider slide using Databricks template.

        Args:
            data: Slide data dict
            prefer_dark: Use dark background (default True for structural slides)
        """
        slide = self._create_slide("section", data, prefer_dark=prefer_dark)

        # Fill title placeholder (idx 0) - template handles styling
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", "Section Title"))

    def add_content_slide(self, data: Dict[str, Any]) -> None:
        """Create content slide with bullets."""
        slide = self._create_slide("content", data)

        # Title
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", "Slide Title"))

        # Subtitle (idx 2 for "7 Content A - Basic" layout)
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=2)
            if subtitle_ph:
                self.fill_text(subtitle_ph, data["subtitle"])

        # Bullets (idx 1)
        body_ph = self.get_placeholder(slide, idx=1)
        self.fill_bullets(body_ph, data.get("bullets", []))

    def add_two_column_slide(self, data: Dict[str, Any]) -> None:
        """Create two-column slide."""
        slide = self._create_slide("two-column", data)

        # Title (type TITLE = 1)
        title_ph = self.get_placeholder(slide, ph_type=1)
        self.fill_text(title_ph, data.get("title", "Two Column"))

        # Subtitle row (idx 5 for "9 Content B - 2 Column" layout)
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=5)
            if subtitle_ph:
                self.fill_text(subtitle_ph, data["subtitle"])

        # Get SUBTITLE placeholders (type 4) sorted by position - these are column headers
        subtitle_phs = self.get_placeholders_by_type(slide, 4)
        # Column headers are at y > 1.5 inches (filter out top subtitle row)
        col_headers = [ph for ph in subtitle_phs if ph.top.inches > 1.5]
        col_headers.sort(key=lambda s: s.left)  # Sort left to right

        # Get BODY placeholders (type 2) sorted by position - these are column content
        body_phs = self.get_placeholders_by_type(slide, 2)
        # Column bodies are at y > 2.5 inches
        col_bodies = [ph for ph in body_phs if ph.top.inches > 2.5]
        col_bodies.sort(key=lambda s: s.left)  # Sort left to right

        # Fill headers
        if data.get("left_header") and len(col_headers) > 0:
            self.fill_text(col_headers[0], data["left_header"])
        if data.get("right_header") and len(col_headers) > 1:
            self.fill_text(col_headers[1], data["right_header"])

        # Fill content
        if len(col_bodies) > 0:
            self.fill_bullets(col_bodies[0], data.get("left", []))
        if len(col_bodies) > 1:
            self.fill_bullets(col_bodies[1], data.get("right", []))

    def add_three_column_slide(self, data: Dict[str, Any]) -> None:
        """Create three-column slide."""
        slide = self._create_slide("three-column", data)

        # Title (type TITLE = 1)
        title_ph = self.get_placeholder(slide, ph_type=1)
        self.fill_text(title_ph, data.get("title", "Three Column"))

        # Subtitle row (idx 7 for "11 Content C - 3 Column" layout)
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=7)
            if subtitle_ph:
                self.fill_text(subtitle_ph, data["subtitle"])

        columns = data.get("columns", [])

        # Get SUBTITLE placeholders (type 4) sorted by position - these are column headers
        # Filter out the top subtitle row (used for page subtitle) by checking y position
        subtitle_phs = self.get_placeholders_by_type(slide, 4)
        # Column headers are typically at y > 2 inches, subtitle row is at y < 2
        col_headers = [ph for ph in subtitle_phs if ph.top.inches > 1.5]
        col_headers.sort(key=lambda s: s.left)  # Sort left to right

        # Get BODY placeholders (type 2) sorted by position - these are column content
        body_phs = self.get_placeholders_by_type(slide, 2)
        # Column bodies are typically at y > 2.5 inches
        col_bodies = [ph for ph in body_phs if ph.top.inches > 2.5]
        col_bodies.sort(key=lambda s: s.left)  # Sort left to right

        # Fill column headers
        for i, col in enumerate(columns[:3]):
            if i < len(col_headers) and col.get("header"):
                self.fill_text(col_headers[i], col["header"])

        # Fill column content
        for i, col in enumerate(columns[:3]):
            if i < len(col_bodies):
                self.fill_bullets(col_bodies[i], col.get("items", []))

    def add_big_number_slide(self, data: Dict[str, Any]) -> None:
        """Create big number/hero stat slide."""
        slide = self._create_slide("big-number", data)

        # Big number in title placeholder (idx 0)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("number", "0"), bold=True, color=COLORS["accent"])

        # Description in body (idx 1)
        body_ph = self.get_placeholder(slide, idx=1)
        description = data.get("text", "")
        if data.get("subtitle"):
            description += f"\n{data['subtitle']}"
        self.fill_text(body_ph, description)

    def add_callout_slide(self, data: Dict[str, Any], prefer_dark: bool = True) -> None:
        """Create callout/bold statement slide using Databricks template.

        Args:
            data: Slide data dict
            prefer_dark: Use dark background (default True for structural slides)
        """
        slide = self._create_slide("callout", data, prefer_dark=prefer_dark)

        # Main text in title placeholder (idx 0) - template handles styling
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("text", "Key message"))

        # Source attribution in subtitle placeholder (idx 1)
        if data.get("source"):
            subtitle_ph = self.get_placeholder(slide, idx=1)
            if subtitle_ph:
                self.fill_text(subtitle_ph, f"— {data['source']}")

    def add_quote_slide(self, data: Dict[str, Any], prefer_dark: bool = True) -> None:
        """Create quote/testimonial slide using Databricks template.

        Args:
            data: Slide data dict
            prefer_dark: Use dark background (default True for structural slides)
        """
        slide = self._create_slide("quote", data, prefer_dark=prefer_dark)

        # Quote text in title placeholder (idx 0) - template handles styling
        quote_text = f'"{data.get("quote", "Quote goes here.")}"'
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, quote_text)

        # Attribution in subtitle placeholder (idx 1)
        if data.get("attribution"):
            subtitle_ph = self.get_placeholder(slide, idx=1)
            if subtitle_ph:
                self.fill_text(subtitle_ph, f"— {data['attribution']}")

    def add_closing_slide(self, data: Dict[str, Any], prefer_dark: bool = True) -> None:
        """Create closing/thank you slide.

        Args:
            data: Slide data dict
            prefer_dark: Use dark background (default True for structural slides)
        """
        slide = self._create_slide("closing", data, prefer_dark=prefer_dark)

        # Colors based on background mode
        text_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]

        # The closing layout has pre-designed graphics (Databricks logo, etc.)
        title = data.get("title", "Thank You")

        # Title text - positioned at top
        self.add_textbox(slide, title, 0.75, 0.8, 11.5, 1.2,
                        font_size=48, bold=True, color=text_color,
                        alignment=PP_ALIGN.CENTER)

    # =========================================================================
    # New Template-Based Slides (8 new types)
    # =========================================================================

    def add_two_column_icons_slide(self, data: Dict[str, Any]) -> None:
        """Create two-column slide with icon spots."""
        slide = self._create_slide("two-column-icons", data)

        # Title (idx 0)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", "Two Column with Icons"))

        # Subtitle row (idx 5)
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=5)
            if subtitle_ph:
                self.fill_text(subtitle_ph, data["subtitle"])

        columns = data.get("columns", [])

        # -- Layout adjustment: place icons and shift columns up --
        # Layout has a 1.5" gap (1.92"–3.42") for icon spots.
        # Place icons there and move headers+bodies up for more body room.
        icon_top = Inches(2.10)
        icon_size = Inches(0.50)
        header_top = Inches(2.72)
        body_top = Inches(3.50)
        body_height = Inches(3.30)

        # Column centers for icon placement (based on body placeholder positions)
        col_centers_in = [0.83 + 4.94 / 2, 7.58 + 4.92 / 2]  # 3.30, 10.04

        for i, col in enumerate(columns[:2]):
            icon_name = col.get("icon")
            if icon_name and i < len(col_centers_in):
                self._place_column_icon(
                    slide, icon_name, col_centers_in[i], icon_top, icon_size
                )

        # Headers (idx 3, 4) — moved up
        for i, col in enumerate(columns[:2]):
            header_ph = self.get_placeholder(slide, idx=3+i)
            if header_ph:
                # Read all inherited values before overriding any
                ol, ow, oh = header_ph.left, header_ph.width, header_ph.height
                header_ph.left, header_ph.top = ol, header_top
                header_ph.width, header_ph.height = ow, oh
                if col.get("header"):
                    self.fill_text(header_ph, col["header"])

        # Content (idx 1, 2) — moved up with more height
        for i, col in enumerate(columns[:2]):
            body_ph = self.get_placeholder(slide, idx=1+i)
            if body_ph:
                ol, ow = body_ph.left, body_ph.width
                body_ph.left, body_ph.top = ol, body_top
                body_ph.width, body_ph.height = ow, body_height
            self.fill_bullets(body_ph, col.get("items", []))

    def add_three_column_icons_slide(self, data: Dict[str, Any]) -> None:
        """Create three-column slide with icon spots."""
        slide = self._create_slide("three-column-icons", data)

        # Title (idx 0)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", "Three Column with Icons"))

        # Subtitle row (idx 7)
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=7)
            if subtitle_ph:
                self.fill_text(subtitle_ph, data["subtitle"])

        columns = data.get("columns", [])

        # -- Layout adjustment: place icons and shift columns up --
        # Layout has a 1.5" gap (1.92"–3.42") for icon spots.
        # Place icons there and move headers+bodies up for more body room.
        icon_top = Inches(2.10)
        icon_size = Inches(0.50)
        header_top = Inches(2.72)
        body_top = Inches(3.50)
        body_height = Inches(3.30)

        # Column centers for icon placement (based on body placeholder positions)
        col_centers_in = [
            0.83 + 3.53 / 2,  # Col 1: 2.595
            4.90 + 3.53 / 2,  # Col 2: 6.665
            8.93 + 3.53 / 2,  # Col 3: 10.695
        ]

        for i, col in enumerate(columns[:3]):
            icon_name = col.get("icon")
            if icon_name and i < len(col_centers_in):
                self._place_column_icon(
                    slide, icon_name, col_centers_in[i], icon_top, icon_size
                )

        # Headers (idx 3, 4, 6) — moved up
        header_indices = [3, 4, 6]
        for i, col in enumerate(columns[:3]):
            header_ph = self.get_placeholder(slide, idx=header_indices[i])
            if header_ph:
                # Read all inherited values before overriding any
                ol, ow, oh = header_ph.left, header_ph.width, header_ph.height
                header_ph.left, header_ph.top = ol, header_top
                header_ph.width, header_ph.height = ow, oh
                if col.get("header"):
                    self.fill_text(header_ph, col["header"])

        # Content (idx 1, 2, 5) — moved up with more height
        body_indices = [1, 2, 5]
        for i, col in enumerate(columns[:3]):
            body_ph = self.get_placeholder(slide, idx=body_indices[i])
            if body_ph:
                ol, ow = body_ph.left, body_ph.width
                body_ph.left, body_ph.top = ol, body_top
                body_ph.width, body_ph.height = ow, body_height
            self.fill_bullets(body_ph, col.get("items", []))

    def add_cards_slide(self, data: Dict[str, Any]) -> None:
        """Create three-column cards slide."""
        slide = self._create_slide("cards", data)

        # Title (idx 0)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", "Cards"))

        # Subtitle row (idx 7)
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=7)
            if subtitle_ph:
                self.fill_text(subtitle_ph, data["subtitle"])

        cards = data.get("cards", [])

        # Card headers (idx 4, 5, 6)
        for i, card in enumerate(cards[:3]):
            header_ph = self.get_placeholder(slide, idx=4+i)
            if header_ph and card.get("header"):
                self.fill_text(header_ph, card["header"])

        # Card content (idx 1, 2, 3)
        for i, card in enumerate(cards[:3]):
            body_ph = self.get_placeholder(slide, idx=1+i)
            if card.get("content"):
                self.fill_text(body_ph, card["content"])
            elif card.get("items"):
                self.fill_bullets(body_ph, card["items"])

    def add_card_right_slide(self, data: Dict[str, Any]) -> None:
        """Create slide with content left, card right."""
        slide = self._create_slide("card-right", data)

        # Title (idx 0)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", "Card Right"))

        # Subtitle (idx 3)
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=3)
            if subtitle_ph:
                self.fill_text(subtitle_ph, data["subtitle"])

        # Left content (idx 1)
        left_ph = self.get_placeholder(slide, idx=1)
        if data.get("content"):
            self.fill_text(left_ph, data["content"])
        elif data.get("bullets"):
            self.fill_bullets(left_ph, data["bullets"])

        # Right card area (idx 2) - for diagrams/images/tables
        right_ph = self.get_placeholder(slide, idx=2)
        if data.get("card_content"):
            self.fill_text(right_ph, data["card_content"])

    def add_card_left_slide(self, data: Dict[str, Any]) -> None:
        """Create slide with card left, content right."""
        slide = self._create_slide("card-left", data)

        # Title (idx 0)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", "Card Left"))

        # Subtitle (idx 3)
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=3)
            if subtitle_ph:
                self.fill_text(subtitle_ph, data["subtitle"])

        # Left card area (idx 2)
        left_ph = self.get_placeholder(slide, idx=2)
        if data.get("card_content"):
            self.fill_text(left_ph, data["card_content"])

        # Right content (idx 1)
        right_ph = self.get_placeholder(slide, idx=1)
        if data.get("content"):
            self.fill_text(right_ph, data["content"])
        elif data.get("bullets"):
            self.fill_bullets(right_ph, data["bullets"])

    def add_card_full_slide(self, data: Dict[str, Any]) -> None:
        """Create slide with full-width card."""
        slide = self._create_slide("card-full", data)

        # Title (idx 0)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", "Full Card"))

        # Subtitle (idx 2)
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=2)
            if subtitle_ph:
                self.fill_text(subtitle_ph, data["subtitle"])

        # Card content (idx 1)
        card_ph = self.get_placeholder(slide, idx=1)
        if data.get("content"):
            self.fill_text(card_ph, data["content"])

    def add_one_column_slide(self, data: Dict[str, Any]) -> None:
        """Create one-column text slide."""
        slide = self._create_slide("one-column", data)

        # Title (idx 0)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", ""))

        # Subtitle row (idx 2 for "7 Content A - Basic" layout)
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=2)
            if subtitle_ph:
                self.fill_text(subtitle_ph, data["subtitle"])

        # Body (idx 1)
        body_ph = self.get_placeholder(slide, idx=1)
        if data.get("content"):
            self.fill_text(body_ph, data["content"])
        elif data.get("bullets"):
            self.fill_bullets(body_ph, data["bullets"])

    def add_section_description_slide(self, data: Dict[str, Any]) -> None:
        """Create section slide with description."""
        slide = self._create_slide("section-description", data)

        # Title (idx 0)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", "Section Title"))

        # Subtitle (idx 1)
        subtitle_ph = self.get_placeholder(slide, idx=1)
        if subtitle_ph:
            self.fill_text(subtitle_ph, data.get("subtitle", ""))

        # Body/description (idx 2)
        body_ph = self.get_placeholder(slide, idx=2)
        if data.get("description"):
            self.fill_text(body_ph, data["description"])
        elif data.get("bullets"):
            self.fill_bullets(body_ph, data["bullets"])

    # =========================================================================
    # Hybrid Slides (template background + custom shapes)
    # =========================================================================

    def add_agenda_slide(self, data: Dict[str, Any], prefer_dark: bool = False) -> None:
        """Create agenda slide with custom hexagon numbers.

        Args:
            data: Slide data dict
            prefer_dark: Use dark background (default False for content slides)
        """
        slide = self._create_slide("agenda", data, prefer_dark=prefer_dark)

        # Colors based on background mode
        title_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]
        item_bg_color = COLORS["text_light"] if prefer_dark else COLORS["light_bg"]

        # Title
        self.add_textbox(slide, data.get("title", "Agenda"),
                        0.83, 0.59, 10.0, 0.8,
                        font_size=36, bold=True, color=title_color)

        # Custom agenda items with hexagons
        items = data.get("items", [])
        start_y = 2.0

        for i, item in enumerate(items, 1):
            y_pos = start_y + (i - 1) * 0.9

            # Hexagon for number
            hex_shape = slide.shapes.add_shape(
                MSO_SHAPE.HEXAGON,
                Inches(0.75), Inches(y_pos),
                Inches(0.6), Inches(0.6)
            )
            hex_shape.fill.solid()
            hex_shape.fill.fore_color.rgb = hex_to_rgb(COLORS["accent"])
            hex_shape.line.fill.background()

            # Number in hexagon (always light text on accent)
            tf = hex_shape.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = str(i)
            p.font.name = FONT_FAMILY
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(COLORS["text_light"])
            p.alignment = PP_ALIGN.CENTER
            tf.anchor = MSO_ANCHOR.MIDDLE

            # Item text with background bar
            item_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1.5), Inches(y_pos),
                Inches(8), Inches(0.6)
            )
            item_bg.fill.solid()
            item_bg.fill.fore_color.rgb = hex_to_rgb(item_bg_color)
            item_bg.line.fill.background()

            # Item text (always dark on light bar)
            self.add_textbox(slide, item, 1.7, y_pos + 0.1, 7.5, 0.5,
                           font_size=20, color=COLORS["text_dark"])

    def add_timeline_slide(self, data: Dict[str, Any], prefer_dark: bool = False) -> None:
        """Create timeline/process slide with steps.

        Args:
            data: Slide data dict
            prefer_dark: Use dark background (default False for content slides)
        """
        slide = self._create_slide("timeline", data, prefer_dark=prefer_dark)

        # Colors based on background mode
        title_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]
        body_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]
        secondary_color = COLORS["text_light"] if prefer_dark else COLORS["text_secondary"]

        # Title
        self.add_textbox(slide, data.get("title", "Timeline"),
                        0.83, 0.59, 10.0, 0.8,
                        font_size=36, bold=True, color=title_color)

        steps = data.get("steps", [])
        num_steps = len(steps)
        if num_steps == 0:
            return

        step_width = 10.5 / num_steps
        start_x = 1.4

        # Connecting line (accent color works on both backgrounds)
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(start_x + 0.3), Inches(3.1),
            Inches(step_width * num_steps - 0.6), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = hex_to_rgb(COLORS["accent"])
        line.line.fill.background()

        for i, step in enumerate(steps):
            x_pos = start_x + (i * step_width)

            # Circle for step number (accent fill, always light text)
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x_pos + step_width/2 - 0.35), Inches(2.75),
                Inches(0.7), Inches(0.7)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = hex_to_rgb(COLORS["accent"])
            circle.line.fill.background()

            # Step number (always light on accent)
            tf = circle.text_frame
            p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.font.name = FONT_FAMILY
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(COLORS["text_light"])
            p.alignment = PP_ALIGN.CENTER
            tf.anchor = MSO_ANCHOR.MIDDLE

            # Step title
            self.add_textbox(slide, step.get("title", f"Step {i+1}"),
                           x_pos, 3.7, step_width, 0.6,
                           font_size=16, bold=True, color=body_color,
                           alignment=PP_ALIGN.CENTER)

            # Step description
            if step.get("description"):
                self.add_textbox(slide, step["description"],
                               x_pos, 4.4, step_width, 1.5,
                               font_size=12, color=secondary_color,
                               alignment=PP_ALIGN.CENTER)

    def add_icon_grid_slide(self, data: Dict[str, Any], prefer_dark: bool = False) -> None:
        """Create icon grid slide for features/capabilities.

        Args:
            data: Slide data dict
            prefer_dark: Use dark background (default False for content slides)
        """
        slide = self._create_slide("icon-grid", data, prefer_dark=prefer_dark)

        # Colors based on background mode
        title_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]
        body_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]
        secondary_color = COLORS["text_light"] if prefer_dark else COLORS["text_secondary"]
        circle_fill = COLORS["text_light"] if prefer_dark else COLORS["light_bg"]

        # Title
        self.add_textbox(slide, data.get("title", "Features"),
                        0.83, 0.59, 10.0, 0.8,
                        font_size=36, bold=True, color=title_color)

        items = data.get("items", data.get("features", []))
        num_items = len(items)
        if num_items == 0:
            return

        # Determine grid layout
        if num_items <= 3:
            cols, rows = num_items, 1
        elif num_items <= 6:
            cols, rows = 3, 2
        else:
            cols, rows = 4, 2

        cell_width = 11 / cols
        cell_height = 2.2
        start_x = 1.2
        start_y = 1.8

        for i, item in enumerate(items[:8]):
            col = i % cols
            row = i // cols
            x_pos = start_x + (col * cell_width)
            y_pos = start_y + (row * (cell_height + 0.5))

            # Circle with accent border
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x_pos + cell_width/2 - 0.5), Inches(y_pos),
                Inches(1), Inches(1)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = hex_to_rgb(circle_fill)
            circle.line.color.rgb = hex_to_rgb(COLORS["accent"])
            circle.line.width = Pt(3)

            # Icon (emoji preferred, falls back to first letter of title)
            icon_raw = item.get("icon")
            if icon_raw:
                icon_text = icon_raw if len(icon_raw) <= 2 else icon_raw[0]
            else:
                icon_text = item.get("title", "?")[0].upper()
            tf = circle.text_frame
            p = tf.paragraphs[0]
            p.text = icon_text
            p.font.name = FONT_FAMILY
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(COLORS["accent"])
            p.alignment = PP_ALIGN.CENTER
            tf.anchor = MSO_ANCHOR.MIDDLE

            # Item title
            self.add_textbox(slide, item.get("title", ""),
                           x_pos, y_pos + 1.1, cell_width, 0.5,
                           font_size=14, bold=True, color=body_color,
                           alignment=PP_ALIGN.CENTER)

            # Item description
            if item.get("description"):
                self.add_textbox(slide, item["description"],
                               x_pos, y_pos + 1.55, cell_width, 0.8,
                               font_size=11, color=secondary_color,
                               alignment=PP_ALIGN.CENTER)

    def add_stat_row_slide(self, data: Dict[str, Any], prefer_dark: bool = False) -> None:
        """Create slide with multiple stats in a row."""
        slide = self._create_slide("stat-row", data, prefer_dark=prefer_dark)

        # Dynamic colors based on background
        title_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]
        label_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]

        # Title
        self.add_textbox(slide, data.get("title", "Key Metrics"),
                        0.83, 0.59, 10.0, 0.8,
                        font_size=36, bold=True, color=title_color)

        stats = data.get("stats", [])
        num_stats = len(stats)
        if num_stats == 0:
            return

        stat_width = 11.5 / num_stats
        start_x = 0.9

        for i, stat in enumerate(stats):
            x_pos = start_x + (i * stat_width)

            # Stat value (accent color works on both backgrounds)
            self.add_textbox(slide, stat.get("value", "0"),
                           x_pos, 2.5, stat_width - 0.3, 1.5,
                           font_size=56, bold=True, color=COLORS["accent"],
                           alignment=PP_ALIGN.CENTER)

            # Stat label
            self.add_textbox(slide, stat.get("label", ""),
                           x_pos, 4.2, stat_width - 0.3, 1.0,
                           font_size=16, bold=True, color=label_color,
                           alignment=PP_ALIGN.CENTER)

            # Divider (except after last)
            if i < num_stats - 1:
                divider = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(x_pos + stat_width - 0.15), Inches(2.7),
                    Inches(0.02), Inches(2.5)
                )
                divider.fill.solid()
                divider.fill.fore_color.rgb = hex_to_rgb(COLORS["divider"])
                divider.line.fill.background()

    def add_pros_cons_slide(self, data: Dict[str, Any], prefer_dark: bool = False) -> None:
        """Create pros/cons comparison slide."""
        slide = self._create_slide("pros-cons", data, prefer_dark=prefer_dark)

        # Dynamic colors based on background
        title_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]
        body_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]

        # Title
        self.add_textbox(slide, data.get("title", "Pros & Cons"),
                        0.83, 0.59, 10.0, 0.8,
                        font_size=36, bold=True, color=title_color)

        # Pros header (green works on both backgrounds)
        self.add_textbox(slide, data.get("pros_header", "Pros"),
                        0.75, 1.6, 5.5, 0.5,
                        font_size=20, bold=True, color=COLORS["green"])

        # Pros items
        pros = data.get("pros", [])
        for i, pro in enumerate(pros):
            self.add_textbox(slide, f"✓  {pro}",
                           0.75, 2.2 + (i * 0.6), 5.5, 0.5,
                           font_size=16, color=body_color)

        # Cons header (red works on both backgrounds)
        self.add_textbox(slide, data.get("cons_header", "Cons"),
                        7.0, 1.6, 5.5, 0.5,
                        font_size=20, bold=True, color=COLORS["red"])

        # Cons items
        cons = data.get("cons", [])
        for i, con in enumerate(cons):
            self.add_textbox(slide, f"✗  {con}",
                           7.0, 2.2 + (i * 0.6), 5.5, 0.5,
                           font_size=16, color=body_color)

    def add_comparison_slide(self, data: Dict[str, Any], prefer_dark: bool = False) -> None:
        """Create VS comparison slide."""
        slide = self._create_slide("comparison", data, prefer_dark=prefer_dark)

        # Dynamic colors based on background
        title_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]
        label_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]
        diamond_bg = COLORS["dark_bg"] if prefer_dark else COLORS["accent"]
        diamond_text = COLORS["text_light"]  # Always white on diamond

        # Title
        self.add_textbox(slide, data.get("title", "Comparison"),
                        0.83, 0.59, 10.0, 0.8,
                        font_size=36, bold=True, color=title_color)

        # VS diamond in center
        diamond = slide.shapes.add_shape(
            MSO_SHAPE.DIAMOND,
            Inches(6.166), Inches(3.25),
            Inches(1), Inches(1)
        )
        diamond.fill.solid()
        diamond.fill.fore_color.rgb = hex_to_rgb(diamond_bg)
        diamond.line.fill.background()

        # VS text
        tf = diamond.text_frame
        p = tf.paragraphs[0]
        p.text = "vs."
        p.font.name = FONT_FAMILY
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(diamond_text)
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        # Left label
        self.add_textbox(slide, data.get("left_label", "Option A"),
                        1.5, 5.0, 4.0, 0.6,
                        font_size=20, bold=True, color=label_color,
                        alignment=PP_ALIGN.CENTER)

        # Right label
        self.add_textbox(slide, data.get("right_label", "Option B"),
                        7.833, 5.0, 4.0, 0.6,
                        font_size=20, bold=True, color=label_color,
                        alignment=PP_ALIGN.CENTER)

    def add_checklist_slide(self, data: Dict[str, Any], prefer_dark: bool = False) -> None:
        """Create checklist slide."""
        slide = self._create_slide("checklist", data, prefer_dark=prefer_dark)

        # Dynamic colors based on background
        title_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]
        body_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]
        unchecked_fill = COLORS["text_light"] if prefer_dark else COLORS["light_bg"]

        # Title
        self.add_textbox(slide, data.get("title", "Checklist"),
                        0.83, 0.59, 10.0, 0.8,
                        font_size=36, bold=True, color=title_color)

        items = data.get("items", [])
        start_y = 1.8

        for i, item in enumerate(items):
            y_pos = start_y + (i * 0.7)

            # Checkbox
            checkbox = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.9), Inches(y_pos),
                Inches(0.35), Inches(0.35)
            )
            checkbox.fill.solid()

            is_checked = item.get("checked", False) if isinstance(item, dict) else False
            item_text = item.get("text", item) if isinstance(item, dict) else item

            if is_checked:
                checkbox.fill.fore_color.rgb = hex_to_rgb(COLORS["accent"])
                # Checkmark (always white on accent background)
                tf = checkbox.text_frame
                p = tf.paragraphs[0]
                p.text = "✓"
                p.font.name = FONT_FAMILY
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = hex_to_rgb(COLORS["text_light"])
                p.alignment = PP_ALIGN.CENTER
                tf.anchor = MSO_ANCHOR.MIDDLE
            else:
                checkbox.fill.fore_color.rgb = hex_to_rgb(unchecked_fill)

            checkbox.line.color.rgb = hex_to_rgb(COLORS["accent"])
            checkbox.line.width = Pt(2)

            # Item text
            self.add_textbox(slide, item_text, 1.5, y_pos, 10.0, 0.4,
                           font_size=16, color=body_color)

    def add_logos_slide(self, data: Dict[str, Any], prefer_dark: bool = False) -> None:
        """Create logo display slide."""
        slide = self._create_slide("logos", data, prefer_dark=prefer_dark)

        # Dynamic colors based on background
        title_color = COLORS["text_light"] if prefer_dark else COLORS["text_dark"]
        subtitle_color = COLORS["text_light"] if prefer_dark else COLORS["text_secondary"]

        # Title
        self.add_textbox(slide, data.get("title", "Our Partners"),
                        0.83, 0.59, 10.0, 0.8,
                        font_size=36, bold=True, color=title_color)

        # Subtitle
        if data.get("subtitle"):
            self.add_textbox(slide, data["subtitle"],
                           0.75, 1.3, 11.0, 0.5,
                           font_size=16, color=subtitle_color,
                           alignment=PP_ALIGN.CENTER)

        logos = data.get("logos", [])
        num_logos = len(logos)

        # Grid layout
        if num_logos <= 4:
            cols, rows = num_logos, 1
        elif num_logos <= 8:
            cols, rows = 4, 2
        else:
            cols, rows = 5, 2

        cell_width = 10 / cols
        cell_height = 1.6
        start_x = 1.7
        start_y = 2.5

        # Logo box colors based on background
        box_fill = COLORS["text_light"] if prefer_dark else COLORS["light_bg"]
        box_text = COLORS["text_secondary"]  # Gray works on white boxes

        for i, logo in enumerate(logos[:10]):
            col = i % cols
            row = i // cols
            x_pos = start_x + (col * cell_width)
            y_pos = start_y + (row * (cell_height + 0.3))

            # Logo placeholder box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x_pos), Inches(y_pos),
                Inches(cell_width - 0.4), Inches(cell_height - 0.3)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = hex_to_rgb(box_fill)
            box.line.color.rgb = hex_to_rgb(COLORS["divider"])
            box.line.width = Pt(1)

            # Company name as placeholder
            logo_name = logo if isinstance(logo, str) else logo.get("name", "Company")
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = logo_name
            p.font.name = FONT_FAMILY
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(box_text)
            p.alignment = PP_ALIGN.CENTER
            tf.anchor = MSO_ANCHOR.MIDDLE

    # =========================================================================
    # Architecture Catalog Import
    # =========================================================================

    def _get_catalog_prs(self):
        """Lazy-load the architecture catalog presentation (cached)."""
        if not hasattr(self, '_catalog_prs'):
            if not CATALOG_PPTX_PATH.exists():
                raise FileNotFoundError(
                    f"Architecture catalog not found: {CATALOG_PPTX_PATH}"
                )
            self._catalog_prs = Presentation(str(CATALOG_PPTX_PATH))
        return self._catalog_prs

    def _get_icon_catalog(self):
        """Lazy-load the icon catalog JSON (cached)."""
        if not hasattr(self, '_icon_catalog'):
            if not ICON_CATALOG_PATH.exists():
                raise FileNotFoundError(
                    f"Icon catalog not found: {ICON_CATALOG_PATH}\n"
                    f"Run: python3 scripts/build-icon-catalog.py"
                )
            with open(ICON_CATALOG_PATH, 'r') as f:
                self._icon_catalog = json.load(f)
        return self._icon_catalog

    def _import_catalog_slide(self, data: Dict[str, Any]) -> None:
        """Import a pre-built architecture diagram from the catalog.

        Creates a BLANK slide in the output presentation, then copies
        the source slide's cSld (common slide data) element to preserve
        all shapes, positioning, z-ordering, groups, and connectors.
        Images are remapped from the source to the target presentation.

        Args:
            data: Slide data dict with required "catalog_slide" (1-based index)
                  and optional "notes" for custom speaker notes.
        """
        from lxml import etree

        catalog_slide_num = data.get("catalog_slide")
        if catalog_slide_num is None:
            print("Warning: architecture slide missing 'catalog_slide', skipping")
            return

        catalog_prs = self._get_catalog_prs()
        slide_index = catalog_slide_num - 1  # Convert to 0-based

        if slide_index < 0 or slide_index >= len(catalog_prs.slides):
            print(f"Warning: catalog_slide {catalog_slide_num} out of range "
                  f"(1-{len(catalog_prs.slides)}), skipping")
            return

        source_slide = catalog_prs.slides[slide_index]

        # Create a BLANK slide in the output presentation
        self.slide_count += 1
        blank_layout = self.layouts.get("BLANK", list(self.layouts.values())[0])
        target_slide = self.prs.slides.add_slide(blank_layout)

        # --- Copy cSld (common slide data) from source to target ---
        source_cSld = source_slide._element.find(
            '{http://schemas.openxmlformats.org/presentationml/2006/main}cSld'
        )
        if source_cSld is None:
            print(f"Warning: No cSld found in catalog slide {catalog_slide_num}")
            return

        new_cSld = deepcopy(source_cSld)

        # --- Resolve scheme colors to explicit RGB values ---
        # The catalog's theme may differ from the output template's theme.
        # Convert all schemeClr references to explicit srgbClr values
        # using the catalog's actual theme colors so they render correctly.
        theme_colors = self._get_source_theme_colors(source_slide)
        if theme_colors:
            self._resolve_scheme_colors(new_cSld, theme_colors)

        # --- Inject explicit white background ---
        # The catalog slides may inherit their background from their layout/master
        # which we're not importing. Set an explicit solid white background.
        nsmap_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        nsmap_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        # Remove any existing bg element
        existing_bg = new_cSld.find(f'{{{nsmap_p}}}bg')
        if existing_bg is None:
            existing_bg = new_cSld.find(f'{{{nsmap_a}}}bg')
        if existing_bg is not None:
            new_cSld.remove(existing_bg)

        # Create explicit white solid fill background
        bg_xml = (
            f'<p:bg xmlns:p="{nsmap_p}" xmlns:a="{nsmap_a}">'
            f'<p:bgPr>'
            f'<a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>'
            f'<a:effectLst/>'
            f'</p:bgPr>'
            f'</p:bg>'
        )
        bg_element = etree.fromstring(bg_xml)
        new_cSld.insert(0, bg_element)

        # --- Remap image relationships ---
        nsmap_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        blips = new_cSld.findall(f'.//{{{nsmap_a}}}blip')

        rId_map = {}  # old_rId -> new_rId
        for blip in blips:
            old_rId = blip.get(f'{{{nsmap_r}}}embed')
            if not old_rId or old_rId in rId_map:
                if old_rId and old_rId in rId_map:
                    blip.set(f'{{{nsmap_r}}}embed', rId_map[old_rId])
                continue

            # Get image blob from source slide
            try:
                source_rel = source_slide.part.rels[old_rId]
                image_part = source_rel.target_part
                image_blob = image_part.blob
                content_type = image_part.content_type
            except (KeyError, AttributeError):
                continue

            # Add image to target slide via the package's image handling
            # This deduplicates by SHA1 hash automatically
            new_rId = self._add_image_to_slide(
                target_slide, image_blob, content_type
            )
            rId_map[old_rId] = new_rId
            blip.set(f'{{{nsmap_r}}}embed', new_rId)

        # --- Apply pre-insertion modifications (XML-level) ---
        # Order: snapshot IDs, remove shapes + fix connectors, replace text,
        #        override title, move shapes/groups, add connectors,
        #        recalculate connector endpoints, verify.
        modifications = data.get("modifications", {})
        has_modifications = bool(modifications)

        # 1. Snapshot shape IDs before any removals (for orphan detection)
        pre_ids = self._collect_shape_ids(new_cSld) if has_modifications else {}

        # 2. Remove shapes, clean up empty containers, fix orphaned connectors
        if modifications.get("remove_shapes"):
            self._apply_shape_removals(new_cSld, modifications["remove_shapes"])
            self._cleanup_empty_containers(new_cSld)
            self._remove_orphaned_connectors(new_cSld, pre_ids)

        # 3. Text replacements
        if modifications.get("text_replacements"):
            modified = self._apply_text_replacements(new_cSld, modifications["text_replacements"])
            self._auto_fit_text(new_cSld, modified_txbodies=modified)

        # 4. Title override
        if data.get("title"):
            self._apply_title_override(new_cSld, data["title"])

        # 5. Shape and group moves
        if modifications.get("move_shapes"):
            self._apply_shape_moves(new_cSld, modifications["move_shapes"])
        if modifications.get("move_groups"):
            self._apply_group_moves(new_cSld, modifications["move_groups"])

        # 6. Add explicit connectors (user-specified)
        if modifications.get("add_connectors"):
            self._apply_connector_additions(
                new_cSld, modifications["add_connectors"]
            )

        # 7. Recalculate all connector endpoints to match current positions
        if has_modifications:
            self._update_connector_endpoints(new_cSld)

        # 8. Verify connector integrity
        if has_modifications:
            issues = self._verify_architecture_modifications(new_cSld)
            for issue in issues:
                print(f"  Warning: {issue}", file=sys.stderr)

        # Replace the target slide's cSld with the copied one
        target_cSld = target_slide._element.find(
            f'{{{nsmap_p}}}cSld'
        )
        if target_cSld is not None:
            target_slide._element.replace(target_cSld, new_cSld)
        else:
            target_slide._element.insert(0, new_cSld)

        # Update python-pptx's cached spTree reference after cSld replacement
        # so that subsequent shape additions go to the correct tree.
        # SlideShapes stores the spTree as _grpSp (inherited from _BaseGroupShapes).
        new_spTree = new_cSld.find(f'{{{nsmap_p}}}spTree')
        if new_spTree is not None:
            target_slide.shapes._grpSp = new_spTree

        # --- Apply post-insertion modifications (live slide object) ---
        if modifications.get("overlays"):
            self._apply_overlays(target_slide, modifications["overlays"])

        # --- Handle speaker notes ---
        custom_notes = data.get("notes")
        if custom_notes:
            target_slide.notes_slide.notes_text_frame.text = custom_notes
        else:
            # Copy notes from the source slide
            try:
                if source_slide.has_notes_slide:
                    source_notes = source_slide.notes_slide.notes_text_frame.text
                    if source_notes:
                        target_slide.notes_slide.notes_text_frame.text = source_notes
            except Exception:
                pass

    def _get_source_theme_colors(self, source_slide):
        """Extract the color scheme from the source slide's theme.

        Navigates slide -> layout -> master -> theme to find the color
        scheme, then returns a dict mapping scheme color names
        (e.g. 'accent1', 'dk1') to hex RGB strings (e.g. '1B5162').
        """
        from lxml import etree

        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        try:
            master = source_slide.slide_layout.slide_master
            for rel in master.part.rels.values():
                if 'theme' in rel.reltype:
                    theme_xml = etree.fromstring(rel.target_part.blob)
                    clrScheme = theme_xml.find(f'.//{{{ns_a}}}clrScheme')
                    if clrScheme is None:
                        continue

                    colors = {}
                    for child in clrScheme:
                        name = etree.QName(child.tag).localname
                        for color_elem in child:
                            ctag = etree.QName(color_elem.tag).localname
                            if ctag == 'srgbClr':
                                colors[name] = color_elem.get('val')
                            elif ctag == 'sysClr':
                                colors[name] = color_elem.get(
                                    'lastClr', color_elem.get('val', '000000')
                                )

                    # OOXML aliases: bg1/bg2/tx1/tx2 map to lt1/lt2/dk1/dk2
                    colors['bg1'] = colors.get('lt1', 'FFFFFF')
                    colors['bg2'] = colors.get('lt2', 'EEEEEE')
                    colors['tx1'] = colors.get('dk1', '000000')
                    colors['tx2'] = colors.get('dk2', '000000')
                    return colors
        except Exception:
            pass
        return {}

    def _apply_color_transforms(self, hex_color, modifiers):
        """Apply OOXML color transform children to a base hex color.

        Handles tint, shade, lumMod, lumOff, satMod, satOff in document
        order per the ECMA-376 spec.  Returns the final hex RGB string.
        """
        import colorsys

        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)

        for mod_name, mod_val in modifiers:
            pct = mod_val / 100000.0

            if mod_name == 'tint':
                # Shift toward white: 100% = no change, 0% = white
                r = int(r * pct + 255 * (1 - pct))
                g = int(g * pct + 255 * (1 - pct))
                b = int(b * pct + 255 * (1 - pct))
            elif mod_name == 'shade':
                # Shift toward black: 100% = no change, 0% = black
                r = int(r * pct)
                g = int(g * pct)
                b = int(b * pct)
            elif mod_name in ('lumMod', 'lumOff', 'satMod', 'satOff'):
                h, l, s = colorsys.rgb_to_hls(r / 255, g / 255, b / 255)
                if mod_name == 'lumMod':
                    l *= pct
                elif mod_name == 'lumOff':
                    l += pct
                elif mod_name == 'satMod':
                    s *= pct
                elif mod_name == 'satOff':
                    s += pct
                l = max(0.0, min(1.0, l))
                s = max(0.0, min(1.0, s))
                r2, g2, b2 = colorsys.hls_to_rgb(h, l, s)
                r, g, b = int(r2 * 255), int(g2 * 255), int(b2 * 255)

        r = max(0, min(255, r))
        g = max(0, min(255, g))
        b = max(0, min(255, b))
        return f'{r:02X}{g:02X}{b:02X}'

    def _resolve_scheme_colors(self, element, color_map):
        """Replace all schemeClr references with explicit srgbClr values.

        This ensures imported catalog slides render with the catalog's
        original colors regardless of the output template's theme.
        """
        from lxml import etree

        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        for scheme_clr in element.findall(f'.//{{{ns_a}}}schemeClr'):
            val = scheme_clr.get('val', '')
            base_hex = color_map.get(val)
            if base_hex is None:
                continue

            # Collect transform modifiers and alpha in document order
            modifiers = []
            alpha_elem = None
            for child in scheme_clr:
                tag = etree.QName(child.tag).localname
                child_val = child.get('val')
                if tag == 'alpha':
                    alpha_elem = deepcopy(child)
                elif child_val is not None:
                    try:
                        modifiers.append((tag, int(child_val)))
                    except ValueError:
                        pass

            # Compute final color
            final_hex = (self._apply_color_transforms(base_hex, modifiers)
                         if modifiers else base_hex)

            # Build replacement srgbClr element
            new_elem = etree.Element(f'{{{ns_a}}}srgbClr')
            new_elem.set('val', final_hex)
            if alpha_elem is not None:
                new_elem.append(alpha_elem)

            # Swap in parent
            parent = scheme_clr.getparent()
            if parent is not None:
                parent.replace(scheme_clr, new_elem)

    # =========================================================================
    # Architecture Slide Modification Methods
    # =========================================================================

    _NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    _NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    _EMU_PER_INCH = 914400

    @staticmethod
    def _normalize_text(text):
        """Normalize Unicode whitespace for comparison."""
        text = unicodedata.normalize('NFKC', text)
        text = re.sub(r'[\u00a0\u200b\u200c\u200d\ufeff]', ' ', text)
        text = re.sub(r'\s+', ' ', text)
        return text.strip()

    def _get_all_txbodies(self, element):
        """Find all txBody elements in an XML tree (including inside groups).

        txBody may be in p: namespace (top-level shapes) or a: namespace.
        """
        results = element.findall(f'.//{{{self._NS_P}}}txBody')
        results += element.findall(f'.//{{{self._NS_A}}}txBody')
        return results

    def _concat_paragraph_text(self, p_elem):
        """Concatenate all run texts in a paragraph element."""
        runs = p_elem.findall(f'{{{self._NS_A}}}r')
        return "".join(
            r.findtext(f'{{{self._NS_A}}}t', default='') for r in runs
        )

    def _get_shape_name_from_sp(self, sp_elem):
        """Get shape name from an sp element's nvSpPr/cNvPr."""
        for ns in [self._NS_P, self._NS_A]:
            cNvPr = sp_elem.find(f'{{{ns}}}nvSpPr/{{{ns}}}cNvPr')
            if cNvPr is not None:
                return cNvPr.get('name', '')
        cNvPr = sp_elem.find(f'.//{{{self._NS_P}}}cNvPr')
        if cNvPr is not None:
            return cNvPr.get('name', '')
        return ''

    def _apply_text_replacements(self, cSld, replacements):
        """Replace text in shapes within a cSld XML tree.

        Each replacement is a dict with:
          - find: text to search for (required)
          - replace: replacement text (required)
          - shape_name: limit to shapes with this name (optional)
          - match_index: 0-based index if find text appears multiple times (optional)

        Preserves formatting of the first matching run.
        Returns a set of modified txBody elements (for auto-fit scoping).
        Uses direct element references (not id()) since lxml proxy objects
        can be garbage-collected and their id() values reused.
        """
        modified_txbodies = set()
        for repl in replacements:
            find_text = self._normalize_text(repl.get("find", ""))
            replace_text = repl.get("replace", "")
            target_shape_name = repl.get("shape_name")
            match_index = repl.get("match_index", 0)

            if not find_text:
                continue

            match_count = 0
            for txBody in self._get_all_txbodies(cSld):
                # Check shape_name filter if specified
                if target_shape_name:
                    sp = txBody.getparent()
                    if sp is not None:
                        name = self._get_shape_name_from_sp(sp)
                        if name != target_shape_name:
                            continue

                for p in txBody.findall(f'{{{self._NS_A}}}p'):
                    runs = p.findall(f'{{{self._NS_A}}}r')
                    if not runs:
                        continue

                    full_text = self._normalize_text(self._concat_paragraph_text(p))
                    if find_text not in full_text:
                        continue

                    if match_count != match_index:
                        match_count += 1
                        continue

                    # Perform replacement: put new text in first run, clear others
                    new_full_text = full_text.replace(find_text, replace_text)
                    t_elem = runs[0].find(f'{{{self._NS_A}}}t')
                    if t_elem is not None:
                        t_elem.text = new_full_text
                    for extra_run in runs[1:]:
                        t_elem = extra_run.find(f'{{{self._NS_A}}}t')
                        if t_elem is not None:
                            t_elem.text = ""

                    modified_txbodies.add(txBody)
                    match_count += 1
                    break  # Move to next replacement
                else:
                    continue
                break  # Found match, move to next replacement

        return modified_txbodies

    def _apply_title_override(self, cSld, new_title):
        """Replace the title text on an imported slide.

        Identifies the title by scoring shapes: largest font + topmost position wins.
        """
        best_score = -1
        best_txBody = None
        best_p = None

        for txBody in self._get_all_txbodies(cSld):
            for p in txBody.findall(f'{{{self._NS_A}}}p'):
                runs = p.findall(f'{{{self._NS_A}}}r')
                if not runs:
                    continue

                text = self._normalize_text(self._concat_paragraph_text(p))
                if not text or len(text) <= 1:
                    continue

                # Score: font size (higher = better) + inverted top position (higher on slide = better)
                font_size = 0
                rPr = runs[0].find(f'{{{self._NS_A}}}rPr')
                if rPr is not None:
                    sz = rPr.get('sz')
                    if sz:
                        font_size = int(sz)

                # Get top position from parent shape
                sp = txBody.getparent()
                top = 99999999
                if sp is not None:
                    # spPr may be in p: or a: namespace
                    for ns in [self._NS_P, self._NS_A]:
                        spPr = sp.find(f'{{{ns}}}spPr')
                        if spPr is not None:
                            xfrm = spPr.find(f'{{{self._NS_A}}}xfrm')
                            if xfrm is not None:
                                off = xfrm.find(f'{{{self._NS_A}}}off')
                                if off is not None:
                                    y = off.get('y')
                                    if y:
                                        top = int(y)
                            break

                # Score combines font size (weight 2) and inversed top position
                score = (font_size * 2) + (99999999 - top)
                if score > best_score:
                    best_score = score
                    best_txBody = txBody
                    best_p = p

        if best_p is not None and best_txBody is not None:
            from lxml import etree
            runs = best_p.findall(f'{{{self._NS_A}}}r')
            if runs:
                t_elem = runs[0].find(f'{{{self._NS_A}}}t')
                if t_elem is not None:
                    t_elem.text = new_title
                for extra_run in runs[1:]:
                    t_elem = extra_run.find(f'{{{self._NS_A}}}t')
                    if t_elem is not None:
                        t_elem.text = ""

            # Enable normAutofit so PowerPoint shrinks title text if needed
            bodyPr = best_txBody.find(f'{{{self._NS_A}}}bodyPr')
            if bodyPr is not None:
                for child in list(bodyPr):
                    if etree.QName(child.tag).localname in ('noAutofit', 'spAutoFit'):
                        bodyPr.remove(child)
                etree.SubElement(bodyPr, f'{{{self._NS_A}}}normAutofit')

    def _auto_fit_text(self, cSld, modified_txbodies=None):
        """Enable PowerPoint's native auto-fit on shapes with replaced text.

        Replaces <a:noAutofit/> with <a:normAutofit/> in the bodyPr element
        of modified shapes. This tells PowerPoint to automatically shrink
        text to fit the shape bounds using its real font metrics, which is
        far more accurate than heuristic character-width estimation.

        Args:
            cSld: The cSld XML element
            modified_txbodies: Set of txBody elements that were modified.
                              If None, all shapes are checked (backward compat).
        """
        from lxml import etree
        ns_a = self._NS_A

        for txBody in self._get_all_txbodies(cSld):
            if modified_txbodies is not None and txBody not in modified_txbodies:
                continue

            bodyPr = txBody.find(f'{{{ns_a}}}bodyPr')
            if bodyPr is None:
                continue

            # Remove existing noAutofit or spAutoFit if present
            for child in list(bodyPr):
                localname = etree.QName(child.tag).localname
                if localname in ('noAutofit', 'spAutoFit'):
                    bodyPr.remove(child)

            # Add normAutofit — PowerPoint will shrink text to fit
            etree.SubElement(bodyPr, f'{{{ns_a}}}normAutofit')

    def _apply_overlays(self, slide, overlays):
        """Add overlay shapes on top of an imported slide.

        Each overlay dict specifies:
          - type: "textbox", "rectangle", "rounded_rectangle", "oval", or "icon"
          - text: text content (optional, not used for icon type)
          - left, top, width, height: position in inches (required; for icon,
            omitting width or height preserves aspect ratio)
          - font_size: in points (default 11)
          - font_color: hex color (default "#FFFFFF")
          - fill_color: hex color for shape fill (optional)
          - border_color: hex color for border (optional)
          - bold: boolean (default False)

        Icon-specific fields:
          - icon: name from icon catalog (e.g., "openai", "kafka")
          - image_path: path to a custom image file (used if icon is not set)
        """
        shape_type_map = {
            "textbox": None,  # Use add_textbox
            "rectangle": MSO_SHAPE.RECTANGLE,
            "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
            "oval": MSO_SHAPE.OVAL,
        }

        for overlay in overlays:
            otype = overlay.get("type", "textbox")

            # --- Icon overlay ---
            if otype == "icon":
                self._apply_icon_overlay(slide, overlay)
                continue

            left = Inches(overlay.get("left", 0))
            top = Inches(overlay.get("top", 0))
            width = Inches(overlay.get("width", 2))
            height = Inches(overlay.get("height", 0.5))
            text = overlay.get("text", "")
            font_size = overlay.get("font_size", 11)
            font_color = overlay.get("font_color", "#FFFFFF")
            fill_color = overlay.get("fill_color")
            border_color = overlay.get("border_color")
            bold = overlay.get("bold", False)

            if otype == "textbox":
                shape = slide.shapes.add_textbox(left, top, width, height)
            else:
                mso_type = shape_type_map.get(otype, MSO_SHAPE.RECTANGLE)
                shape = slide.shapes.add_shape(mso_type, left, top, width, height)

            # Apply fill
            if fill_color:
                shape.fill.solid()
                shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
            elif otype != "textbox":
                shape.fill.background()

            # Apply border
            if border_color:
                shape.line.color.rgb = hex_to_rgb(border_color)
                shape.line.width = Pt(1)
            elif otype != "textbox":
                shape.line.fill.background()

            # Add text
            if text:
                tf = shape.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = text
                p.font.name = FONT_FAMILY
                p.font.size = Pt(font_size)
                p.font.bold = bold
                p.font.color.rgb = hex_to_rgb(font_color)
                p.alignment = PP_ALIGN.CENTER

    def _place_column_icon(self, slide, icon_name: str,
                           col_center_in: float, icon_top, icon_size) -> None:
        """Place an icon centered in a column's icon spot.

        Args:
            slide: The slide to add the icon to
            icon_name: Icon name from the icon catalog
            col_center_in: Horizontal center of the column in inches
            icon_top: Top position (Emu)
            icon_size: Width and height (Emu, square)
        """
        try:
            catalog = self._get_icon_catalog()
            icon_entry = catalog.get("icons", {}).get(icon_name)
            if not icon_entry:
                print(f"Warning: icon '{icon_name}' not found in catalog")
                return
            icon_path = str(ICON_DIR / icon_entry["file"])
            icon_left = Inches(col_center_in) - icon_size // 2
            slide.shapes.add_picture(icon_path, icon_left, icon_top,
                                     icon_size, icon_size)
        except Exception as e:
            print(f"Warning: could not place icon '{icon_name}': {e}")

    def _apply_icon_overlay(self, slide, overlay):
        """Add an icon image overlay on top of a slide.

        Resolves the image from either:
          - icon: name from the icon catalog (e.g., "openai")
          - image_path: direct path to an image file

        Position/size via left, top, width, height (inches).
        If only width or height is given, the other is set to None
        so python-pptx preserves the aspect ratio.
        """
        icon_name = overlay.get("icon")
        image_path = overlay.get("image_path")

        if icon_name:
            catalog = self._get_icon_catalog()
            icon_entry = catalog.get("icons", {}).get(icon_name)
            if not icon_entry:
                print(f"Warning: icon '{icon_name}' not found in catalog, skipping")
                return
            resolved_path = str(ICON_DIR / icon_entry["file"])
        elif image_path:
            resolved_path = str(Path(image_path).expanduser())
            if not Path(resolved_path).exists():
                print(f"Warning: image_path '{image_path}' not found, skipping")
                return
        else:
            print("Warning: icon overlay missing both 'icon' and 'image_path', skipping")
            return

        left = Inches(overlay.get("left", 0))
        top = Inches(overlay.get("top", 0))

        # If only one dimension is given, set the other to None
        # so python-pptx preserves the original aspect ratio
        has_width = "width" in overlay
        has_height = "height" in overlay
        width = Inches(overlay["width"]) if has_width else None
        height = Inches(overlay["height"]) if has_height else None

        slide.shapes.add_picture(resolved_path, left, top, width, height)

    def _apply_shape_removals(self, cSld, removals):
        """Remove shapes from a cSld XML tree by text content or shape name.

        Each removal dict has:
          - text: match shapes containing this text (optional)
          - shape_name: match shapes by name attribute (optional)
          - remove_group: if true (default), remove the entire parent group
            when the matched shape is inside a group. This is the typical
            behavior for reference architecture source/consumer labels where
            the icon and label are grouped together.

        At least one of text/shape_name must be provided.
        Works for both top-level and grouped shapes.
        """
        ns_p = self._NS_P
        ns_a = self._NS_A

        for removal in removals:
            find_text = removal.get("text")
            find_name = removal.get("shape_name")
            remove_group = removal.get("remove_group", True)

            if not find_text and not find_name:
                continue

            if find_text:
                find_text = self._normalize_text(find_text)

            # Find all sp elements (including in groups)
            for sp in cSld.findall(f'.//{{{ns_p}}}sp'):
                matched = False

                if find_name:
                    name = self._get_shape_name_from_sp(sp)
                    if name == find_name:
                        matched = True

                if find_text and not matched:
                    txBody = sp.find(f'{{{ns_p}}}txBody')
                    if txBody is None:
                        txBody = sp.find(f'{{{ns_a}}}txBody')
                    if txBody is not None:
                        full = ""
                        for p in txBody.findall(f'{{{ns_a}}}p'):
                            full += self._concat_paragraph_text(p) + " "
                        full = self._normalize_text(full)
                        if find_text in full:
                            matched = True

                if matched:
                    parent = sp.getparent()
                    if parent is None:
                        continue

                    # If shape is inside a group and remove_group is enabled,
                    # remove the entire parent group (e.g. icon + label unit)
                    if remove_group and parent.tag == f'{{{ns_p}}}grpSp':
                        group_parent = parent.getparent()
                        if group_parent is not None:
                            group_parent.remove(parent)
                    else:
                        parent.remove(sp)

                        # Check if parent group is now empty
                        if parent.tag == f'{{{ns_p}}}grpSp':
                            remaining = parent.findall(f'{{{ns_p}}}sp')
                            remaining += parent.findall(f'{{{ns_p}}}grpSp')
                            remaining += parent.findall(f'{{{ns_p}}}cxnSp')
                            remaining += parent.findall(f'{{{ns_p}}}pic')
                            if not remaining:
                                gp = parent.getparent()
                                if gp is not None:
                                    gp.remove(parent)
                    break  # One removal per entry

    def _cleanup_empty_containers(self, cSld):
        """After shape removals, detect and remove empty container boxes.

        Finds large AUTO_SHAPE rectangles with no text content (bounding boxes),
        checks if any content groups still overlap them spatially, and removes
        empty ones. Adjacent non-empty containers are expanded to fill the
        vacated vertical space.

        Also removes orphaned labels (small text shapes) whose vertical center
        falls within the bounds of a removed container.
        """
        from lxml import etree

        ns_p = self._NS_P
        ns_a = self._NS_A

        spTree = cSld.find(f'{{{ns_p}}}spTree')
        if spTree is None:
            return

        EMU_HALF_INCH = 457200
        EMU_TENTH_INCH = 91440

        def get_sp_xfrm(sp_elem):
            """Get (x, y, cx, cy) in EMU from a shape's spPr/xfrm."""
            for ns in [ns_p, ns_a]:
                spPr = sp_elem.find(f'{{{ns}}}spPr')
                if spPr is not None:
                    xfrm = spPr.find(f'{{{ns_a}}}xfrm')
                    if xfrm is not None:
                        off = xfrm.find(f'{{{ns_a}}}off')
                        ext = xfrm.find(f'{{{ns_a}}}ext')
                        if off is not None and ext is not None:
                            return (
                                int(off.get('x', '0')),
                                int(off.get('y', '0')),
                                int(ext.get('cx', '0')),
                                int(ext.get('cy', '0')),
                                off, ext,
                            )
            return None

        def set_sp_pos(off_elem, ext_elem, y, cy):
            """Update a shape's y position and height."""
            off_elem.set('y', str(int(y)))
            ext_elem.set('cy', str(int(cy)))

        def get_grp_xfrm(grp_elem):
            """Get (x, y, cx, cy, off_elem) from a group's grpSpPr/xfrm."""
            for ns in [ns_p, ns_a]:
                grpSpPr = grp_elem.find(f'{{{ns}}}grpSpPr')
                if grpSpPr is not None:
                    xfrm = grpSpPr.find(f'{{{ns_a}}}xfrm')
                    if xfrm is not None:
                        off = xfrm.find(f'{{{ns_a}}}off')
                        ext = xfrm.find(f'{{{ns_a}}}ext')
                        if off is not None and ext is not None:
                            return (
                                int(off.get('x', '0')),
                                int(off.get('y', '0')),
                                int(ext.get('cx', '0')),
                                int(ext.get('cy', '0')),
                                off,
                            )
            return None

        def get_sp_text(sp_elem):
            """Get concatenated text from a shape."""
            text = ""
            for txBody in [sp_elem.find(f'{{{ns_p}}}txBody'),
                           sp_elem.find(f'{{{ns_a}}}txBody')]:
                if txBody is not None:
                    for p in txBody.findall(f'{{{ns_a}}}p'):
                        text += self._concat_paragraph_text(p) + " "
            return text.strip()

        def overlaps(ax, ay, acx, acy, bx, by, bcx, bcy):
            """Check if two rectangles overlap spatially."""
            return (ax < bx + bcx and ax + acx > bx and
                    ay < by + bcy and ay + acy > by)

        # --- Collect shapes by type ---
        containers = []   # (elem, x, y, cx, cy, off, ext)
        content_groups = []  # (x, y, cx, cy, off_elem)
        label_shapes = []    # (elem, x, y, cx, cy, text)

        for elem in list(spTree):
            tag = etree.QName(elem.tag).localname

            if tag == 'sp':
                xfrm_data = get_sp_xfrm(elem)
                if xfrm_data is None:
                    continue
                x, y, cx, cy, off, ext = xfrm_data
                text = get_sp_text(elem)

                if not text and cx > EMU_HALF_INCH and cy > EMU_HALF_INCH:
                    # Large shape with no text = potential container
                    containers.append((elem, x, y, cx, cy, off, ext))
                elif text and cy < EMU_HALF_INCH:
                    # Small text shape = potential label
                    label_shapes.append((elem, x, y, cx, cy, text))

            elif tag == 'grpSp':
                grp_data = get_grp_xfrm(elem)
                if grp_data is not None:
                    content_groups.append(grp_data)

        if not containers:
            return

        # --- Classify containers as empty or non-empty ---
        empty = []
        non_empty = []

        for cdata in containers:
            elem, cx_, cy_, ccx, ccy, off, ext = cdata
            has_content = False
            for gx, gy, gcx, gcy, _ in content_groups:
                if overlaps(cx_, cy_, ccx, ccy, gx, gy, gcx, gcy):
                    has_content = True
                    break
            if has_content:
                non_empty.append(cdata)
            else:
                empty.append(cdata)

        if not empty:
            return

        # --- Remove empty containers and orphaned labels ---
        for elem, ex, ey, ecx, ecy, _, _ in empty:
            try:
                spTree.remove(elem)
            except ValueError:
                pass

            # Remove labels whose vertical center falls within the empty box
            for ldata in label_shapes[:]:
                lelem, lx, ly, lcx, lcy, ltext = ldata
                label_cy_mid = ly + lcy // 2
                if ey <= label_cy_mid <= ey + ecy and lx < ex + ecx:
                    try:
                        spTree.remove(lelem)
                        label_shapes.remove(ldata)
                    except ValueError:
                        pass

        # --- Expand adjacent non-empty containers to fill gaps ---
        # Group containers by column (similar x position), only in columns
        # that had empty containers removed.
        def same_column(x1, x2):
            return abs(x1 - x2) < EMU_TENTH_INCH * 2

        # Start from empty container columns only
        columns = {}  # column_x -> [(y, cy, is_empty, off_elem, ext_elem)]
        for elem, x, y, cx, cy, off, ext in empty:
            col_key = None
            for k in columns:
                if same_column(k, x):
                    col_key = k
                    break
            if col_key is None:
                col_key = x
                columns[col_key] = []
            columns[col_key].append((y, cy, True, None, None))

        # Only add non-empty containers that share a column with an empty one
        for elem, x, y, cx, cy, off, ext in non_empty:
            col_key = None
            for k in columns:
                if same_column(k, x):
                    col_key = k
                    break
            if col_key is not None:
                columns[col_key].append((y, cy, False, off, ext))

        expanded_containers = []  # track which containers were actually expanded

        for col_key, items in columns.items():
            items.sort(key=lambda t: t[0])  # sort by y position

            remaining = [(y, cy, off, ext) for y, cy, is_empty, off, ext
                         in items if not is_empty]
            if not remaining:
                continue

            # Get the full vertical range of the column
            all_tops = [y for y, cy, _, _, _ in items]
            all_bottoms = [y + cy for y, cy, _, _, _ in items]
            col_top = min(all_tops)
            col_bottom = max(all_bottoms)

            if len(remaining) == 1:
                # Single remaining box — expand to cover full column range
                y, cy, off, ext = remaining[0]
                set_sp_pos(off, ext, col_top, col_bottom - col_top)
                expanded_containers.append((
                    int(off.get('x', '0')), col_top,
                    int(ext.get('cx', '0')), col_bottom - col_top))
            else:
                # Multiple remaining — expand first upward, last downward
                # First box: expand up to column top
                y0, cy0, off0, ext0 = remaining[0]
                if y0 > col_top:
                    set_sp_pos(off0, ext0, col_top, (y0 + cy0) - col_top)
                expanded_containers.append((
                    int(off0.get('x', '0')), int(off0.get('y', '0')),
                    int(ext0.get('cx', '0')), int(ext0.get('cy', '0'))))

                # Last box: expand down to column bottom
                yn, cyn, offn, extn = remaining[-1]
                if yn + cyn < col_bottom:
                    set_sp_pos(offn, extn, yn, col_bottom - yn)
                expanded_containers.append((
                    int(offn.get('x', '0')), int(offn.get('y', '0')),
                    int(extn.get('cx', '0')), int(extn.get('cy', '0'))))

                # Middle boxes: expand to fill gaps between neighbors
                for i in range(len(remaining) - 1):
                    yi, cyi, offi, exti = remaining[i]
                    yi1, _, _, _ = remaining[i + 1]
                    gap = yi1 - (yi + cyi)
                    if gap > EMU_TENTH_INCH:
                        set_sp_pos(offi, exti, yi, yi1 - yi)
                    if i > 0:
                        expanded_containers.append((
                            int(offi.get('x', '0')), int(offi.get('y', '0')),
                            int(exti.get('cx', '0')), int(exti.get('cy', '0'))))

        for ldata in label_shapes:
            lelem, lx, ly, lcx, lcy, ltext = ldata
            lxfrm = get_sp_xfrm(lelem)
            if lxfrm is None:
                continue
            _, _, _, _, l_off, l_ext = lxfrm

            # Find the container this label belongs to:
            # label is to the left of or overlapping the container horizontally,
            # and label's vertical center is within the container's bounds
            label_vmid = ly + lcy // 2
            best_container = None
            for cx_, cy_, ccx, ccy in expanded_containers:
                if (lx < cx_ + ccx and
                        cy_ <= label_vmid <= cy_ + ccy):
                    best_container = (cx_, cy_, ccx, ccy)
                    break

            if best_container is not None:
                _, cont_y, _, cont_cy = best_container
                # Center label vertically within container
                new_label_y = cont_y + (cont_cy - lcy) // 2
                l_off.set('y', str(int(new_label_y)))

        # --- Reposition content groups to be centered in their container ---
        for cx_, cy_, ccx, ccy in expanded_containers:
            # Find all content groups inside this container
            contained = []
            for gx, gy, gcx, gcy, g_off in content_groups:
                if overlaps(cx_, cy_, ccx, ccy, gx, gy, gcx, gcy):
                    contained.append((gx, gy, gcx, gcy, g_off))

            if not contained:
                continue

            # Calculate the bounding box of all contained groups
            content_top = min(gy for _, gy, _, _, _ in contained)
            content_bottom = max(gy + gcy for _, gy, _, gcy, _ in contained)
            content_height = content_bottom - content_top

            # Calculate the shift to center the content block in the container
            container_center = cy_ + ccy // 2
            content_center = content_top + content_height // 2
            shift_y = container_center - content_center

            if abs(shift_y) < EMU_TENTH_INCH:
                continue  # Already centered

            # Apply the shift to each group
            for _, gy, _, _, g_off in contained:
                new_gy = gy + shift_y
                g_off.set('y', str(int(new_gy)))

    def _get_group_transform(self, grpSp_elem):
        """Extract group coordinate transform from a grpSp element.

        Returns a dict with off_x/y, ext_cx/cy, chOff_x/y, chExt_cx/cy,
        has_rotation, and _xfrm_elem (for in-place expansion), or None.
        """
        ns_a = self._NS_A
        ns_p = self._NS_P

        for ns in [ns_p, ns_a]:
            grpSpPr = grpSp_elem.find(f'{{{ns}}}grpSpPr')
            if grpSpPr is not None:
                xfrm = grpSpPr.find(f'{{{ns_a}}}xfrm')
                if xfrm is not None:
                    off = xfrm.find(f'{{{ns_a}}}off')
                    ext = xfrm.find(f'{{{ns_a}}}ext')
                    chOff = xfrm.find(f'{{{ns_a}}}chOff')
                    chExt = xfrm.find(f'{{{ns_a}}}chExt')
                    if all(e is not None for e in [off, ext, chOff, chExt]):
                        return {
                            "off_x": int(off.get('x', '0')),
                            "off_y": int(off.get('y', '0')),
                            "ext_cx": int(ext.get('cx', '0')),
                            "ext_cy": int(ext.get('cy', '0')),
                            "chOff_x": int(chOff.get('x', '0')),
                            "chOff_y": int(chOff.get('y', '0')),
                            "chExt_cx": int(chExt.get('cx', '0')),
                            "chExt_cy": int(chExt.get('cy', '0')),
                            "has_rotation": xfrm.get('rot') is not None,
                            "_off": off,
                            "_ext": ext,
                            "_chOff": chOff,
                            "_chExt": chExt,
                        }
        return None

    def _slide_to_child_coords(self, slide_x_emu, slide_y_emu, gxfrm):
        """Convert slide-space EMU coordinates to child-space EMU.

        Uses: child = chOff + (slide - off) * (chExt / ext)
        """
        if gxfrm["ext_cx"] == 0 or gxfrm["ext_cy"] == 0:
            return slide_x_emu, slide_y_emu

        child_x = gxfrm["chOff_x"] + int(
            (slide_x_emu - gxfrm["off_x"])
            * gxfrm["chExt_cx"] / gxfrm["ext_cx"]
        )
        child_y = gxfrm["chOff_y"] + int(
            (slide_y_emu - gxfrm["off_y"])
            * gxfrm["chExt_cy"] / gxfrm["ext_cy"]
        )
        return child_x, child_y

    def _child_to_slide_coords(self, child_x_emu, child_y_emu, gxfrm):
        """Convert child-space EMU coordinates to slide-space EMU.

        Uses: slide = off + (child - chOff) * (ext / chExt)
        """
        if gxfrm["chExt_cx"] == 0 or gxfrm["chExt_cy"] == 0:
            return child_x_emu, child_y_emu

        slide_x = gxfrm["off_x"] + int(
            (child_x_emu - gxfrm["chOff_x"])
            * gxfrm["ext_cx"] / gxfrm["chExt_cx"]
        )
        slide_y = gxfrm["off_y"] + int(
            (child_y_emu - gxfrm["chOff_y"])
            * gxfrm["ext_cy"] / gxfrm["chExt_cy"]
        )
        return slide_x, slide_y

    def _expand_group_bounds(self, gxfrm, child_x, child_y, child_w, child_h):
        """Expand group bounds in-place if child shape exceeds current area.

        Adjusts off/ext/chOff/chExt so that:
        - The child shape fits within the group's child coordinate space
        - Existing children don't visually shift (parent off shifts proportionally)
        """
        ch_right = child_x + child_w
        ch_bottom = child_y + child_h

        old_chOff_x = gxfrm["chOff_x"]
        old_chOff_y = gxfrm["chOff_y"]
        old_chExt_cx = gxfrm["chExt_cx"]
        old_chExt_cy = gxfrm["chExt_cy"]

        new_chOff_x = min(old_chOff_x, child_x)
        new_chOff_y = min(old_chOff_y, child_y)
        new_chRight = max(old_chOff_x + old_chExt_cx, ch_right)
        new_chBottom = max(old_chOff_y + old_chExt_cy, ch_bottom)
        new_chExt_cx = new_chRight - new_chOff_x
        new_chExt_cy = new_chBottom - new_chOff_y

        if (new_chOff_x == old_chOff_x and new_chOff_y == old_chOff_y
                and new_chExt_cx == old_chExt_cx and new_chExt_cy == old_chExt_cy):
            return  # No expansion needed

        # Compute scale factors (child-to-parent)
        sx = gxfrm["ext_cx"] / old_chExt_cx if old_chExt_cx else 1.0
        sy = gxfrm["ext_cy"] / old_chExt_cy if old_chExt_cy else 1.0

        # Shift parent off proportionally when chOff moves left/up
        delta_chOff_x = new_chOff_x - old_chOff_x  # negative = leftward
        delta_chOff_y = new_chOff_y - old_chOff_y  # negative = upward
        new_off_x = gxfrm["off_x"] + int(delta_chOff_x * sx)
        new_off_y = gxfrm["off_y"] + int(delta_chOff_y * sy)

        # Expand parent ext proportionally
        new_ext_cx = int(new_chExt_cx * sx)
        new_ext_cy = int(new_chExt_cy * sy)

        # Update XML elements in-place
        gxfrm["_off"].set('x', str(new_off_x))
        gxfrm["_off"].set('y', str(new_off_y))
        gxfrm["_ext"].set('cx', str(new_ext_cx))
        gxfrm["_ext"].set('cy', str(new_ext_cy))
        gxfrm["_chOff"].set('x', str(new_chOff_x))
        gxfrm["_chOff"].set('y', str(new_chOff_y))
        gxfrm["_chExt"].set('cx', str(new_chExt_cx))
        gxfrm["_chExt"].set('cy', str(new_chExt_cy))

        # Update cached values so subsequent calls see the new bounds
        gxfrm["off_x"] = new_off_x
        gxfrm["off_y"] = new_off_y
        gxfrm["ext_cx"] = new_ext_cx
        gxfrm["ext_cy"] = new_ext_cy
        gxfrm["chOff_x"] = new_chOff_x
        gxfrm["chOff_y"] = new_chOff_y
        gxfrm["chExt_cx"] = new_chExt_cx
        gxfrm["chExt_cy"] = new_chExt_cy

    def _apply_shape_moves(self, cSld, moves):
        """Reposition shapes within a cSld XML tree.

        Each move dict has:
          - text: find shape by text content (optional)
          - shape_name: find shape by name (optional)
          - left: new x position in inches (optional)
          - top: new y position in inches (optional)

        Supports both top-level and grouped shapes. For grouped shapes,
        slide-space coordinates are converted to the group's child-space.
        Rotated groups emit a warning (coordinates may be inaccurate).
        """
        ns_p = self._NS_P
        ns_a = self._NS_A

        for move in moves:
            find_text = move.get("text")
            find_name = move.get("shape_name")
            new_left = move.get("left")
            new_top = move.get("top")

            if not find_text and not find_name:
                continue
            if new_left is None and new_top is None:
                continue

            if find_text:
                find_text = self._normalize_text(find_text)

            for sp in cSld.findall(f'.//{{{ns_p}}}sp'):
                matched = False

                if find_name:
                    name = self._get_shape_name_from_sp(sp)
                    if name == find_name:
                        matched = True

                if find_text and not matched:
                    txBody = sp.find(f'{{{ns_p}}}txBody')
                    if txBody is None:
                        txBody = sp.find(f'{{{ns_a}}}txBody')
                    if txBody is not None:
                        full = ""
                        for p in txBody.findall(f'{{{ns_a}}}p'):
                            full += self._concat_paragraph_text(p) + " "
                        full = self._normalize_text(full)
                        if find_text in full:
                            matched = True

                if not matched:
                    continue

                # Find shape's spPr and current offset
                spPr = None
                for spPr_ns in [ns_p, ns_a]:
                    spPr = sp.find(f'{{{spPr_ns}}}spPr')
                    if spPr is not None:
                        break
                if spPr is None:
                    break

                xfrm = spPr.find(f'{{{ns_a}}}xfrm')
                if xfrm is None:
                    break
                off = xfrm.find(f'{{{ns_a}}}off')
                ext = xfrm.find(f'{{{ns_a}}}ext')
                if off is None:
                    break

                parent = sp.getparent()
                is_grouped = (parent is not None
                              and parent.tag == f'{{{ns_p}}}grpSp')

                if is_grouped:
                    gxfrm = self._get_group_transform(parent)
                    if gxfrm is None:
                        break

                    if gxfrm["has_rotation"]:
                        print(f"Warning: moving shape in rotated group — "
                              f"position may be inaccurate", file=sys.stderr)

                    # Current child-space position
                    cur_child_x = int(off.get('x', '0'))
                    cur_child_y = int(off.get('y', '0'))

                    # For unmoved axes, back-compute slide-space from child-space
                    cur_slide_x, cur_slide_y = self._child_to_slide_coords(
                        cur_child_x, cur_child_y, gxfrm
                    )

                    target_x_emu = (int(new_left * self._EMU_PER_INCH)
                                    if new_left is not None
                                    else cur_slide_x)
                    target_y_emu = (int(new_top * self._EMU_PER_INCH)
                                    if new_top is not None
                                    else cur_slide_y)

                    # Get child shape dimensions for bounds expansion
                    child_w = int(ext.get('cx', '0')) if ext is not None else 0
                    child_h = int(ext.get('cy', '0')) if ext is not None else 0

                    # Convert target to child-space
                    child_x, child_y = self._slide_to_child_coords(
                        target_x_emu, target_y_emu, gxfrm
                    )

                    # Expand group bounds if shape would exceed them
                    self._expand_group_bounds(
                        gxfrm, child_x, child_y, child_w, child_h
                    )

                    # Recalculate child coords after expansion (chOff may have shifted)
                    child_x, child_y = self._slide_to_child_coords(
                        target_x_emu, target_y_emu, gxfrm
                    )

                    off.set('x', str(child_x))
                    off.set('y', str(child_y))
                else:
                    # Top-level shape: direct EMU assignment
                    if new_left is not None:
                        off.set('x', str(int(new_left * self._EMU_PER_INCH)))
                    if new_top is not None:
                        off.set('y', str(int(new_top * self._EMU_PER_INCH)))

                break  # One move per entry

    def _apply_group_moves(self, cSld, moves):
        """Reposition entire groups within a cSld XML tree.

        Each move dict has:
          - text: find group containing a child shape with this text (optional)
          - group_name: find group by its own name (optional)
          - left: new x position in inches (optional)
          - top: new y position in inches (optional)

        Moves all children together by updating the group's own offset.
        """
        ns_p = self._NS_P
        ns_a = self._NS_A

        for move in moves:
            find_text = move.get("text")
            find_group_name = move.get("group_name")
            new_left = move.get("left")
            new_top = move.get("top")

            if not find_text and not find_group_name:
                continue
            if new_left is None and new_top is None:
                continue

            if find_text:
                find_text = self._normalize_text(find_text)

            for grpSp in cSld.findall(f'.//{{{ns_p}}}grpSp'):
                matched = False

                if find_group_name:
                    # Check group's own name in nvGrpSpPr/cNvPr
                    for ns in [ns_p, ns_a]:
                        cNvPr = grpSp.find(f'{{{ns}}}nvGrpSpPr/{{{ns}}}cNvPr')
                        if cNvPr is not None:
                            if cNvPr.get('name', '') == find_group_name:
                                matched = True
                            break

                if find_text and not matched:
                    # Search child shapes for matching text
                    for sp in grpSp.findall(f'{{{ns_p}}}sp'):
                        txBody = sp.find(f'{{{ns_p}}}txBody')
                        if txBody is None:
                            txBody = sp.find(f'{{{ns_a}}}txBody')
                        if txBody is not None:
                            full = ""
                            for p in txBody.findall(f'{{{ns_a}}}p'):
                                full += self._concat_paragraph_text(p) + " "
                            full = self._normalize_text(full)
                            if find_text in full:
                                matched = True
                                break

                if not matched:
                    continue

                # Update the group's own offset
                gxfrm = self._get_group_transform(grpSp)
                if gxfrm is None:
                    break

                if new_left is not None:
                    gxfrm["_off"].set('x', str(int(new_left * self._EMU_PER_INCH)))
                if new_top is not None:
                    gxfrm["_off"].set('y', str(int(new_top * self._EMU_PER_INCH)))

                break  # One move per entry

    # =========================================================================
    # Connector Management Methods
    # =========================================================================

    def _collect_shape_ids(self, cSld) -> Dict[str, Any]:
        """Build a mapping of cNvPr id -> element for all shapes in a cSld.

        Covers sp, grpSp, and pic elements at all nesting levels.
        Connector shapes (cxnSp) are excluded since they reference
        these IDs rather than being referenced by them.

        Returns:
            Dict mapping id string -> XML element.
        """
        ns_p = self._NS_P
        ns_a = self._NS_A
        id_map: Dict[str, Any] = {}

        for sp in cSld.findall(f'.//{{{ns_p}}}sp'):
            for ns in [ns_p, ns_a]:
                cNvPr = sp.find(f'{{{ns}}}nvSpPr/{{{ns}}}cNvPr')
                if cNvPr is not None:
                    shape_id = cNvPr.get('id')
                    if shape_id:
                        id_map[shape_id] = sp
                    break

        for grp in cSld.findall(f'.//{{{ns_p}}}grpSp'):
            for ns in [ns_p, ns_a]:
                cNvPr = grp.find(f'{{{ns}}}nvGrpSpPr/{{{ns}}}cNvPr')
                if cNvPr is not None:
                    shape_id = cNvPr.get('id')
                    if shape_id:
                        id_map[shape_id] = grp
                    break

        for pic in cSld.findall(f'.//{{{ns_p}}}pic'):
            for ns in [ns_p, ns_a]:
                cNvPr = pic.find(f'{{{ns}}}nvPicPr/{{{ns}}}cNvPr')
                if cNvPr is not None:
                    shape_id = cNvPr.get('id')
                    if shape_id:
                        id_map[shape_id] = pic
                    break

        return id_map

    def _get_max_shape_id(self, cSld) -> int:
        """Find the highest cNvPr id value in the cSld tree.

        Used to generate unique IDs for new connector shapes.
        """
        ns_p = self._NS_P
        max_id = 0
        for cNvPr in cSld.findall(f'.//{{{ns_p}}}cNvPr'):
            try:
                val = int(cNvPr.get('id', '0'))
                if val > max_id:
                    max_id = val
            except ValueError:
                pass
        return max_id

    def _get_shape_bounds_in_slide_coords(
        self, shape_elem
    ) -> Optional[Tuple[int, int, int, int]]:
        """Get a shape's bounding box in slide-space EMU coordinates.

        Handles top-level shapes and shapes nested inside one or more
        groups by walking up the parent chain and applying group
        coordinate transforms at each level.

        Args:
            shape_elem: An sp, grpSp, or pic XML element.

        Returns:
            (x, y, cx, cy) in EMU slide coordinates, or None if
            the shape's transform could not be resolved.
        """
        from lxml import etree

        ns_a = self._NS_A
        ns_p = self._NS_P

        tag = etree.QName(shape_elem.tag).localname
        x: Optional[int] = None
        y: Optional[int] = None
        cx: Optional[int] = None
        cy: Optional[int] = None

        if tag == 'grpSp':
            gxfrm = self._get_group_transform(shape_elem)
            if gxfrm is None:
                return None
            x, y = gxfrm["off_x"], gxfrm["off_y"]
            cx, cy = gxfrm["ext_cx"], gxfrm["ext_cy"]
        else:
            for ns in [ns_p, ns_a]:
                spPr = shape_elem.find(f'{{{ns}}}spPr')
                if spPr is not None:
                    xfrm = spPr.find(f'{{{ns_a}}}xfrm')
                    if xfrm is not None:
                        off = xfrm.find(f'{{{ns_a}}}off')
                        ext = xfrm.find(f'{{{ns_a}}}ext')
                        if off is not None and ext is not None:
                            x = int(off.get('x', '0'))
                            y = int(off.get('y', '0'))
                            cx = int(ext.get('cx', '0'))
                            cy = int(ext.get('cy', '0'))
                            break

        if x is None:
            return None

        # Walk up parent chain, converting through group transforms
        current = shape_elem
        while True:
            parent = current.getparent()
            if parent is None:
                break
            parent_tag = etree.QName(parent.tag).localname
            if parent_tag != 'grpSp':
                break

            gxfrm = self._get_group_transform(parent)
            if gxfrm is None:
                break

            if gxfrm["chExt_cx"] > 0 and gxfrm["chExt_cy"] > 0:
                sx = gxfrm["ext_cx"] / gxfrm["chExt_cx"]
                sy = gxfrm["ext_cy"] / gxfrm["chExt_cy"]

                x = gxfrm["off_x"] + int((x - gxfrm["chOff_x"]) * sx)
                y = gxfrm["off_y"] + int((y - gxfrm["chOff_y"]) * sy)
                cx = int(cx * sx)
                cy = int(cy * sy)

            current = parent

        return (x, y, cx, cy)

    @staticmethod
    def _get_connection_point(x: int, y: int, cx: int, cy: int,
                              idx: int) -> Tuple[int, int]:
        """Calculate a connection point position from shape bounds.

        Standard OOXML connection point indices for rectangles:
          0 = top center
          1 = right center
          2 = bottom center
          3 = left center

        Args:
            x, y, cx, cy: Shape bounding box in EMU.
            idx: Connection point index.

        Returns:
            (point_x, point_y) in EMU.
        """
        points = {
            0: (x + cx // 2, y),
            1: (x + cx, y + cy // 2),
            2: (x + cx // 2, y + cy),
            3: (x, y + cy // 2),
        }
        return points.get(idx, (x + cx // 2, y + cy // 2))

    @staticmethod
    def _get_named_connection_index(name: str) -> int:
        """Convert a named connection point to an index.

        Accepts: "top", "right", "bottom", "left".
        Returns -1 for unrecognized names.
        """
        name_map = {"top": 0, "right": 1, "bottom": 2, "left": 3}
        return name_map.get(name.lower(), -1)

    def _get_all_connectors(self, cSld) -> List:
        """Find all cxnSp elements in a cSld tree."""
        return cSld.findall(f'.//{{{self._NS_P}}}cxnSp')

    def _get_connector_refs(
        self, cxnSp
    ) -> Tuple[Optional[Tuple[str, int]], Optional[Tuple[str, int]]]:
        """Extract start and end connection references from a connector.

        Returns:
            Tuple of (start_ref, end_ref) where each is either
            (shape_id_str, connection_point_idx) or None.
        """
        ns_p = self._NS_P
        ns_a = self._NS_A

        start_ref: Optional[Tuple[str, int]] = None
        end_ref: Optional[Tuple[str, int]] = None

        for ns in [ns_p, ns_a]:
            cNvCxnSpPr = cxnSp.find(
                f'{{{ns}}}nvCxnSpPr/{{{ns}}}cNvCxnSpPr'
            )
            if cNvCxnSpPr is not None:
                stCxn = cNvCxnSpPr.find(f'{{{ns_a}}}stCxn')
                if stCxn is not None:
                    sid = stCxn.get('id')
                    sidx = int(stCxn.get('idx', '0'))
                    if sid:
                        start_ref = (sid, sidx)

                endCxn = cNvCxnSpPr.find(f'{{{ns_a}}}endCxn')
                if endCxn is not None:
                    eid = endCxn.get('id')
                    eidx = int(endCxn.get('idx', '0'))
                    if eid:
                        end_ref = (eid, eidx)
                break

        return start_ref, end_ref

    def _get_connector_line_style(self, cxnSp) -> Dict[str, Any]:
        """Extract line style properties from a connector element.

        Returns dict with keys: color, width, dash, head_end, tail_end.
        """
        ns_p = self._NS_P
        ns_a = self._NS_A
        style: Dict[str, Any] = {}

        spPr = cxnSp.find(f'{{{ns_p}}}spPr')
        if spPr is None:
            return style

        ln = spPr.find(f'{{{ns_a}}}ln')
        if ln is None:
            return style

        w = ln.get('w')
        if w:
            style['width'] = int(w)

        solidFill = ln.find(f'{{{ns_a}}}solidFill')
        if solidFill is not None:
            srgbClr = solidFill.find(f'{{{ns_a}}}srgbClr')
            if srgbClr is not None:
                style['color'] = srgbClr.get('val', '000000')

        prstDash = ln.find(f'{{{ns_a}}}prstDash')
        if prstDash is not None:
            style['dash'] = prstDash.get('val')

        tailEnd = ln.find(f'{{{ns_a}}}tailEnd')
        if tailEnd is not None:
            style['tail_end'] = tailEnd.get('type', 'none')

        headEnd = ln.find(f'{{{ns_a}}}headEnd')
        if headEnd is not None:
            style['head_end'] = headEnd.get('type', 'none')

        return style

    def _build_connector_element(
        self,
        shape_id: int,
        name: str,
        start_pt: Tuple[int, int],
        end_pt: Tuple[int, int],
        start_shape_id: Optional[str] = None,
        start_idx: int = 0,
        end_shape_id: Optional[str] = None,
        end_idx: int = 0,
        line_color: str = '000000',
        line_width: int = 12700,
        arrow: str = 'end',
        dash_style: Optional[str] = None,
    ) -> Any:
        """Build a cxnSp XML element for a straight connector.

        Args:
            shape_id: Unique shape ID for the new connector.
            name: Display name for the connector.
            start_pt: (x, y) start position in EMU.
            end_pt: (x, y) end position in EMU.
            start_shape_id: ID of shape at start (for stCxn reference).
            start_idx: Connection point index at start shape.
            end_shape_id: ID of shape at end (for endCxn reference).
            end_idx: Connection point index at end shape.
            line_color: Hex color string without '#' (default black).
            line_width: Line width in EMU (default ~1pt = 12700).
            arrow: Arrow style — 'end', 'start', 'both', 'none'.
            dash_style: Dash preset (e.g. 'dash', 'dot') or None for solid.

        Returns:
            lxml Element for the connector shape.
        """
        from lxml import etree

        ns_p = self._NS_P
        ns_a = self._NS_A

        sx, sy = start_pt
        ex, ey = end_pt

        x = min(sx, ex)
        y = min(sy, ey)
        cx_val = abs(ex - sx)
        cy_val = abs(ey - sy)
        flip_h = '1' if sx > ex else None
        flip_v = '1' if sy > ey else None

        cxnSp = etree.Element(f'{{{ns_p}}}cxnSp')

        # nvCxnSpPr
        nvCxnSpPr = etree.SubElement(cxnSp, f'{{{ns_p}}}nvCxnSpPr')
        cNvPr = etree.SubElement(nvCxnSpPr, f'{{{ns_p}}}cNvPr')
        cNvPr.set('id', str(shape_id))
        cNvPr.set('name', name)

        cNvCxnSpPr = etree.SubElement(nvCxnSpPr, f'{{{ns_p}}}cNvCxnSpPr')
        if start_shape_id is not None:
            stCxn = etree.SubElement(cNvCxnSpPr, f'{{{ns_a}}}stCxn')
            stCxn.set('id', str(start_shape_id))
            stCxn.set('idx', str(start_idx))
        if end_shape_id is not None:
            endCxn = etree.SubElement(cNvCxnSpPr, f'{{{ns_a}}}endCxn')
            endCxn.set('id', str(end_shape_id))
            endCxn.set('idx', str(end_idx))

        etree.SubElement(nvCxnSpPr, f'{{{ns_p}}}nvPr')

        # spPr
        spPr = etree.SubElement(cxnSp, f'{{{ns_p}}}spPr')
        xfrm = etree.SubElement(spPr, f'{{{ns_a}}}xfrm')
        if flip_h:
            xfrm.set('flipH', flip_h)
        if flip_v:
            xfrm.set('flipV', flip_v)

        off = etree.SubElement(xfrm, f'{{{ns_a}}}off')
        off.set('x', str(x))
        off.set('y', str(y))
        ext = etree.SubElement(xfrm, f'{{{ns_a}}}ext')
        ext.set('cx', str(cx_val))
        ext.set('cy', str(cy_val))

        prstGeom = etree.SubElement(spPr, f'{{{ns_a}}}prstGeom')
        prstGeom.set('prst', 'straightConnector1')
        etree.SubElement(prstGeom, f'{{{ns_a}}}avLst')

        ln = etree.SubElement(spPr, f'{{{ns_a}}}ln')
        ln.set('w', str(line_width))
        solidFill = etree.SubElement(ln, f'{{{ns_a}}}solidFill')
        srgbClr = etree.SubElement(solidFill, f'{{{ns_a}}}srgbClr')
        srgbClr.set('val', line_color.lstrip('#'))

        if dash_style:
            prstDash = etree.SubElement(ln, f'{{{ns_a}}}prstDash')
            prstDash.set('val', dash_style)

        if arrow in ('end', 'both'):
            tailEnd = etree.SubElement(ln, f'{{{ns_a}}}tailEnd')
            tailEnd.set('type', 'triangle')
        if arrow in ('start', 'both'):
            headEnd = etree.SubElement(ln, f'{{{ns_a}}}headEnd')
            headEnd.set('type', 'triangle')

        return cxnSp

    def _remove_orphaned_connectors(
        self, cSld, pre_removal_ids: Dict[str, Any]
    ) -> None:
        """Remove connectors that reference deleted shapes, with auto-reconnection.

        When an intermediate shape is removed and it had exactly one incoming
        and one outgoing connector, a bridge connector is created to maintain
        the logical flow (A -> B -> C becomes A -> C when B is removed).

        Connectors with both endpoints orphaned are removed entirely.
        Connectors with one orphaned endpoint have the reference removed.

        Args:
            cSld: The cSld XML element after shape removals.
            pre_removal_ids: Shape ID map from before removals
                             (from _collect_shape_ids).
        """
        ns_p = self._NS_P
        ns_a = self._NS_A

        current_ids = self._collect_shape_ids(cSld)
        removed_ids = set(pre_removal_ids.keys()) - set(current_ids.keys())

        if not removed_ids:
            return

        connectors = self._get_all_connectors(cSld)
        if not connectors:
            return

        # Build connectivity graph for auto-reconnection
        # Maps: shape_id -> list of (peer_id, own_idx, peer_idx, connector, style)
        incoming: Dict[str, List] = {}
        outgoing: Dict[str, List] = {}

        for cxn in connectors:
            start_ref, end_ref = self._get_connector_refs(cxn)
            if start_ref and end_ref:
                from_id, from_idx = start_ref
                to_id, to_idx = end_ref
                style = self._get_connector_line_style(cxn)

                outgoing.setdefault(from_id, []).append(
                    (to_id, from_idx, to_idx, cxn, style)
                )
                incoming.setdefault(to_id, []).append(
                    (from_id, from_idx, to_idx, cxn, style)
                )

        # Auto-reconnect: bridge linear chains through removed nodes
        next_id = self._get_max_shape_id(cSld) + 1
        bridge_count = 0
        spTree = cSld.find(f'{{{ns_p}}}spTree')

        for removed_id in removed_ids:
            inc = incoming.get(removed_id, [])
            out = outgoing.get(removed_id, [])

            # Only auto-reconnect clear linear chains (1 in, 1 out)
            if len(inc) != 1 or len(out) != 1:
                continue

            upstream_id, _, inc_to_idx, _, inc_style = inc[0]
            downstream_id, out_from_idx, _, _, out_style = out[0]

            # Skip if either neighbor was also removed
            if upstream_id in removed_ids or downstream_id in removed_ids:
                continue

            upstream_elem = current_ids.get(upstream_id)
            downstream_elem = current_ids.get(downstream_id)
            if upstream_elem is None or downstream_elem is None:
                continue

            up_bounds = self._get_shape_bounds_in_slide_coords(upstream_elem)
            down_bounds = self._get_shape_bounds_in_slide_coords(
                downstream_elem
            )
            if up_bounds is None or down_bounds is None:
                continue

            start_pt = self._get_connection_point(*up_bounds, out_from_idx)
            end_pt = self._get_connection_point(*down_bounds, inc_to_idx)

            bridge_color = out_style.get(
                'color', inc_style.get('color', '000000')
            )
            bridge_width = out_style.get(
                'width', inc_style.get('width', 12700)
            )
            bridge_dash = out_style.get('dash', inc_style.get('dash'))

            bridge = self._build_connector_element(
                shape_id=next_id,
                name=f'Bridge Connector {bridge_count + 1}',
                start_pt=start_pt,
                end_pt=end_pt,
                start_shape_id=upstream_id,
                start_idx=out_from_idx,
                end_shape_id=downstream_id,
                end_idx=inc_to_idx,
                line_color=bridge_color,
                line_width=bridge_width,
                arrow='end',
                dash_style=bridge_dash,
            )

            if spTree is not None:
                spTree.append(bridge)
                next_id += 1
                bridge_count += 1

        if bridge_count:
            print(f"  Auto-reconnected: {bridge_count} bridge connector(s)")

        # Remove orphaned connectors
        removed_count = 0
        for cxn in connectors:
            start_ref, end_ref = self._get_connector_refs(cxn)
            start_orphaned = (
                start_ref is not None and start_ref[0] in removed_ids
            )
            end_orphaned = (
                end_ref is not None and end_ref[0] in removed_ids
            )

            if start_orphaned and end_orphaned:
                parent = cxn.getparent()
                if parent is not None:
                    parent.remove(cxn)
                    removed_count += 1
            elif start_orphaned or end_orphaned:
                for ns in [ns_p, ns_a]:
                    cNvCxnSpPr = cxn.find(
                        f'{{{ns}}}nvCxnSpPr/{{{ns}}}cNvCxnSpPr'
                    )
                    if cNvCxnSpPr is not None:
                        if start_orphaned:
                            stCxn = cNvCxnSpPr.find(f'{{{ns_a}}}stCxn')
                            if stCxn is not None:
                                cNvCxnSpPr.remove(stCxn)
                        if end_orphaned:
                            endCxn = cNvCxnSpPr.find(f'{{{ns_a}}}endCxn')
                            if endCxn is not None:
                                cNvCxnSpPr.remove(endCxn)
                        break

                # Remove if no references remain at all
                remaining_start, remaining_end = self._get_connector_refs(cxn)
                if remaining_start is None and remaining_end is None:
                    parent = cxn.getparent()
                    if parent is not None:
                        parent.remove(cxn)
                        removed_count += 1

        if removed_count:
            print(f"  Removed: {removed_count} orphaned connector(s)")

    def _find_shape_by_text_or_name(
        self,
        cSld,
        text: Optional[str] = None,
        shape_name: Optional[str] = None,
    ) -> Optional[Tuple[str, Any]]:
        """Find a shape by text content or name, returning (id, element).

        Searches sp elements first, then groups (matching by child text).

        Args:
            cSld: The cSld XML element to search.
            text: Text to search for in shapes (partial match).
            shape_name: Exact shape name to match.

        Returns:
            (shape_id, shape_element) or None if not found.
        """
        ns_p = self._NS_P
        ns_a = self._NS_A

        if text:
            text = self._normalize_text(text)

        for sp in cSld.findall(f'.//{{{ns_p}}}sp'):
            if shape_name:
                name = self._get_shape_name_from_sp(sp)
                if name == shape_name:
                    for ns in [ns_p, ns_a]:
                        cNvPr = sp.find(f'{{{ns}}}nvSpPr/{{{ns}}}cNvPr')
                        if cNvPr is not None:
                            return (cNvPr.get('id'), sp)

            if text:
                txBody = sp.find(f'{{{ns_p}}}txBody')
                if txBody is None:
                    txBody = sp.find(f'{{{ns_a}}}txBody')
                if txBody is not None:
                    full = ""
                    for p in txBody.findall(f'{{{ns_a}}}p'):
                        full += self._concat_paragraph_text(p) + " "
                    full = self._normalize_text(full)
                    if text in full:
                        for ns in [ns_p, ns_a]:
                            cNvPr = sp.find(
                                f'{{{ns}}}nvSpPr/{{{ns}}}cNvPr'
                            )
                            if cNvPr is not None:
                                return (cNvPr.get('id'), sp)

        # Also search groups (match by child text or group name)
        for grp in cSld.findall(f'.//{{{ns_p}}}grpSp'):
            if shape_name:
                for ns in [ns_p, ns_a]:
                    cNvPr = grp.find(
                        f'{{{ns}}}nvGrpSpPr/{{{ns}}}cNvPr'
                    )
                    if cNvPr is not None:
                        if cNvPr.get('name', '') == shape_name:
                            return (cNvPr.get('id'), grp)
                        break

            if text:
                for sp in grp.findall(f'{{{ns_p}}}sp'):
                    txBody = sp.find(f'{{{ns_p}}}txBody')
                    if txBody is None:
                        txBody = sp.find(f'{{{ns_a}}}txBody')
                    if txBody is not None:
                        full = ""
                        for p in txBody.findall(f'{{{ns_a}}}p'):
                            full += self._concat_paragraph_text(p) + " "
                        full = self._normalize_text(full)
                        if text in full:
                            for ns in [ns_p, ns_a]:
                                cNvPr = grp.find(
                                    f'{{{ns}}}nvGrpSpPr/{{{ns}}}cNvPr'
                                )
                                if cNvPr is not None:
                                    return (cNvPr.get('id'), grp)
                            break

        return None

    def _auto_detect_connection_points(
        self,
        from_bounds: Tuple[int, int, int, int],
        to_bounds: Tuple[int, int, int, int],
    ) -> Tuple[int, int]:
        """Auto-detect the best connection point indices based on position.

        Compares shape centers and picks the most natural connection
        (right->left for horizontal, bottom->top for vertical).

        Returns:
            (from_idx, to_idx) where idx is 0=top, 1=right, 2=bottom, 3=left.
        """
        fx, fy, fcx, fcy = from_bounds
        tx, ty, tcx, tcy = to_bounds

        from_cx = fx + fcx // 2
        from_cy = fy + fcy // 2
        to_cx = tx + tcx // 2
        to_cy = ty + tcy // 2

        dx = to_cx - from_cx
        dy = to_cy - from_cy

        if abs(dx) > abs(dy):
            if dx > 0:
                return (1, 3)  # right -> left
            else:
                return (3, 1)  # left -> right
        else:
            if dy > 0:
                return (2, 0)  # bottom -> top
            else:
                return (0, 2)  # top -> bottom

    def _apply_connector_additions(
        self, cSld, connectors: List[Dict[str, Any]]
    ) -> None:
        """Add new connector shapes between specified components.

        Each connector dict has:
          - from_text or from_name: identify the source shape
          - to_text or to_name: identify the target shape
          - from_point: "top", "right", "bottom", "left" (optional,
                        auto-detected if omitted)
          - to_point: same as from_point (optional)
          - line_color: hex color string (optional, default "#000000")
          - line_width: width in points (optional, default 1)
          - dash_style: "dash", "dot", etc. or omit for solid (optional)
          - arrow: "end", "start", "both", "none" (optional, default "end")
        """
        ns_p = self._NS_P
        spTree = cSld.find(f'{{{ns_p}}}spTree')
        if spTree is None:
            return

        next_id = self._get_max_shape_id(cSld) + 1
        added_count = 0

        for conn_spec in connectors:
            from_text = conn_spec.get("from_text")
            from_name = conn_spec.get("from_name")
            to_text = conn_spec.get("to_text")
            to_name = conn_spec.get("to_name")

            from_result = self._find_shape_by_text_or_name(
                cSld, text=from_text, shape_name=from_name
            )
            to_result = self._find_shape_by_text_or_name(
                cSld, text=to_text, shape_name=to_name
            )

            if from_result is None:
                label = from_text or from_name
                print(
                    f"Warning: add_connectors — source shape "
                    f"'{label}' not found, skipping",
                    file=sys.stderr,
                )
                continue
            if to_result is None:
                label = to_text or to_name
                print(
                    f"Warning: add_connectors — target shape "
                    f"'{label}' not found, skipping",
                    file=sys.stderr,
                )
                continue

            from_id, from_elem = from_result
            to_id, to_elem = to_result

            from_bounds = self._get_shape_bounds_in_slide_coords(from_elem)
            to_bounds = self._get_shape_bounds_in_slide_coords(to_elem)
            if from_bounds is None or to_bounds is None:
                print(
                    "Warning: add_connectors — could not resolve shape "
                    "bounds, skipping",
                    file=sys.stderr,
                )
                continue

            # Determine connection points
            from_point_name = conn_spec.get("from_point")
            to_point_name = conn_spec.get("to_point")

            if from_point_name and to_point_name:
                from_idx = self._get_named_connection_index(from_point_name)
                to_idx = self._get_named_connection_index(to_point_name)
                if from_idx == -1:
                    from_idx = 1
                if to_idx == -1:
                    to_idx = 3
            else:
                from_idx, to_idx = self._auto_detect_connection_points(
                    from_bounds, to_bounds
                )

            start_pt = self._get_connection_point(*from_bounds, from_idx)
            end_pt = self._get_connection_point(*to_bounds, to_idx)

            line_color = conn_spec.get("line_color", "#000000").lstrip('#')
            line_width_pt = conn_spec.get("line_width", 1.0)
            line_width_emu = int(line_width_pt * 12700)
            arrow = conn_spec.get("arrow", "end")
            dash = conn_spec.get("dash_style")
            if dash == "solid":
                dash = None

            cxn_elem = self._build_connector_element(
                shape_id=next_id,
                name=f'Added Connector {added_count + 1}',
                start_pt=start_pt,
                end_pt=end_pt,
                start_shape_id=from_id,
                start_idx=from_idx,
                end_shape_id=to_id,
                end_idx=to_idx,
                line_color=line_color,
                line_width=line_width_emu,
                arrow=arrow,
                dash_style=dash,
            )

            spTree.append(cxn_elem)
            next_id += 1
            added_count += 1

        if added_count:
            print(f"  Added: {added_count} new connector(s)")

    def _update_connector_endpoints(self, cSld) -> None:
        """Recalculate xfrm for all connectors based on their connected shapes.

        For each connector with stCxn and/or endCxn references, looks up the
        current position of the referenced shapes and recalculates the
        connector's xfrm so the arrow endpoints match the shapes' current
        positions. This is critical for Google Slides which uses xfrm
        directly rather than recalculating from connection references.

        Should be called after all shape moves, group moves, and container
        expansion are complete.
        """
        ns_p = self._NS_P
        ns_a = self._NS_A

        current_ids = self._collect_shape_ids(cSld)
        updated_count = 0

        for cxn in self._get_all_connectors(cSld):
            start_ref, end_ref = self._get_connector_refs(cxn)

            if start_ref is None and end_ref is None:
                continue

            spPr = cxn.find(f'{{{ns_p}}}spPr')
            if spPr is None:
                continue
            xfrm = spPr.find(f'{{{ns_a}}}xfrm')
            if xfrm is None:
                continue
            off = xfrm.find(f'{{{ns_a}}}off')
            ext = xfrm.find(f'{{{ns_a}}}ext')
            if off is None or ext is None:
                continue

            cur_x = int(off.get('x', '0'))
            cur_y = int(off.get('y', '0'))
            cur_cx = int(ext.get('cx', '0'))
            cur_cy = int(ext.get('cy', '0'))
            flip_h = xfrm.get('flipH') == '1'
            flip_v = xfrm.get('flipV') == '1'

            # Determine current start and end points from xfrm + flips
            if not flip_h and not flip_v:
                old_start = (cur_x, cur_y)
                old_end = (cur_x + cur_cx, cur_y + cur_cy)
            elif flip_h and not flip_v:
                old_start = (cur_x + cur_cx, cur_y)
                old_end = (cur_x, cur_y + cur_cy)
            elif not flip_h and flip_v:
                old_start = (cur_x, cur_y + cur_cy)
                old_end = (cur_x + cur_cx, cur_y)
            else:
                old_start = (cur_x + cur_cx, cur_y + cur_cy)
                old_end = (cur_x, cur_y)

            new_start = old_start
            new_end = old_end

            if start_ref:
                s_id, s_idx = start_ref
                s_elem = current_ids.get(s_id)
                if s_elem is not None:
                    bounds = self._get_shape_bounds_in_slide_coords(s_elem)
                    if bounds is not None:
                        new_start = self._get_connection_point(*bounds, s_idx)

            if end_ref:
                e_id, e_idx = end_ref
                e_elem = current_ids.get(e_id)
                if e_elem is not None:
                    bounds = self._get_shape_bounds_in_slide_coords(e_elem)
                    if bounds is not None:
                        new_end = self._get_connection_point(*bounds, e_idx)

            if new_start == old_start and new_end == old_end:
                continue

            # Recalculate xfrm from new endpoints
            sx, sy = new_start
            ex, ey = new_end
            new_x = min(sx, ex)
            new_y = min(sy, ey)
            new_cx = abs(ex - sx)
            new_cy = abs(ey - sy)
            new_flip_h = sx > ex
            new_flip_v = sy > ey

            off.set('x', str(new_x))
            off.set('y', str(new_y))
            ext.set('cx', str(new_cx))
            ext.set('cy', str(new_cy))

            if new_flip_h:
                xfrm.set('flipH', '1')
            elif 'flipH' in xfrm.attrib:
                del xfrm.attrib['flipH']

            if new_flip_v:
                xfrm.set('flipV', '1')
            elif 'flipV' in xfrm.attrib:
                del xfrm.attrib['flipV']

            updated_count += 1

        if updated_count:
            print(f"  Updated: {updated_count} connector endpoint(s)")

    def _verify_architecture_modifications(self, cSld) -> List[str]:
        """Validate connector integrity after all modifications.

        Checks:
          - All stCxn/endCxn references point to existing shapes
          - Connector xfrm positions are within slide bounds
          - Reports diagnostic summary

        Returns:
            List of warning strings (empty if all checks pass).
        """
        ns_p = self._NS_P
        ns_a = self._NS_A
        issues: List[str] = []

        current_ids = self._collect_shape_ids(cSld)
        connectors = self._get_all_connectors(cSld)

        # Slide bounds in EMU (13.333 x 7.5 inches)
        slide_width = int(13.333 * self._EMU_PER_INCH)
        slide_height = int(7.5 * self._EMU_PER_INCH)
        margin = self._EMU_PER_INCH  # 1 inch tolerance

        connected_count = 0
        freestanding_count = 0
        orphan_ref_count = 0

        for cxn in connectors:
            start_ref, end_ref = self._get_connector_refs(cxn)

            if start_ref is None and end_ref is None:
                freestanding_count += 1
            else:
                connected_count += 1

            if start_ref and start_ref[0] not in current_ids:
                orphan_ref_count += 1
                issues.append(
                    f"Connector references non-existent start shape "
                    f"id={start_ref[0]}"
                )

            if end_ref and end_ref[0] not in current_ids:
                orphan_ref_count += 1
                issues.append(
                    f"Connector references non-existent end shape "
                    f"id={end_ref[0]}"
                )

            spPr = cxn.find(f'{{{ns_p}}}spPr')
            if spPr is not None:
                xfrm = spPr.find(f'{{{ns_a}}}xfrm')
                if xfrm is not None:
                    off_elem = xfrm.find(f'{{{ns_a}}}off')
                    ext_elem = xfrm.find(f'{{{ns_a}}}ext')
                    if off_elem is not None and ext_elem is not None:
                        cx_pos = int(off_elem.get('x', '0'))
                        cy_pos = int(off_elem.get('y', '0'))
                        cw = int(ext_elem.get('cx', '0'))
                        ch = int(ext_elem.get('cy', '0'))

                        if (cx_pos + cw < -margin
                                or cx_pos > slide_width + margin
                                or cy_pos + ch < -margin
                                or cy_pos > slide_height + margin):
                            issues.append(
                                f"Connector at ({cx_pos}, {cy_pos}) "
                                f"size ({cw}x{ch}) is off-slide"
                            )

        total_shapes = len(current_ids)
        total_connectors = len(connectors)
        print(
            f"  Verification: {total_shapes} shapes, "
            f"{total_connectors} connectors "
            f"({connected_count} connected, "
            f"{freestanding_count} free-standing)"
        )
        if orphan_ref_count:
            print(
                f"  WARNING: {orphan_ref_count} orphaned connector "
                f"reference(s) found"
            )
        if not issues:
            print(f"  All connector references valid")

        return issues

    def _add_image_to_slide(self, slide, image_blob: bytes,
                            content_type: str) -> str:
        """Add an image blob to a slide and return its relationship ID.

        Uses python-pptx's ImagePart for SHA1-based deduplication.
        """
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        from pptx.parts.image import ImagePart

        # Determine file extension from content type
        ext_map = {
            'image/png': '.png',
            'image/jpeg': '.jpg',
            'image/gif': '.gif',
            'image/svg+xml': '.svg',
            'image/x-emf': '.emf',
            'image/x-wmf': '.wmf',
            'image/tiff': '.tiff',
            'image/bmp': '.bmp',
        }
        ext = ext_map.get(content_type, '.png')

        # Check if this exact image already exists in the package (by blob hash)
        package = slide.part.package
        for part in package.iter_parts():
            if (part.content_type == content_type and
                    hasattr(part, 'blob') and part.blob == image_blob):
                # Image already in package, just add relationship to this slide
                return slide.part.relate_to(part, RT.IMAGE)

        # Create new image part
        partname = self._next_image_partname(package, ext)
        image_part = ImagePart(partname, content_type, package, image_blob)
        return slide.part.relate_to(image_part, RT.IMAGE)

    def _next_image_partname(self, package, ext: str):
        """Generate the next available image part name."""
        from pptx.opc.packuri import PackURI

        existing = set()
        for part in package.iter_parts():
            existing.add(str(part.partname))

        n = 1
        while True:
            candidate = f'/ppt/media/image{n}{ext}'
            if candidate not in existing:
                return PackURI(candidate)
            n += 1

    # =========================================================================
    # Generation Methods
    # =========================================================================

    def generate(self, content: Dict[str, Any]) -> None:
        """Generate all slides from content dict."""
        slides = content.get("slides", [])

        slide_methods = {
            # Direct template slides
            "title": self.add_title_slide,
            "section": self.add_section_slide,
            "content": self.add_content_slide,
            "two-column": self.add_two_column_slide,
            "three-column": self.add_three_column_slide,
            "big-number": self.add_big_number_slide,
            "callout": self.add_callout_slide,
            "quote": self.add_quote_slide,
            "closing": self.add_closing_slide,
            # New template slides
            "two-column-icons": self.add_two_column_icons_slide,
            "three-column-icons": self.add_three_column_icons_slide,
            "cards": self.add_cards_slide,
            "card-right": self.add_card_right_slide,
            "card-left": self.add_card_left_slide,
            "card-full": self.add_card_full_slide,
            "one-column": self.add_one_column_slide,
            "section-description": self.add_section_description_slide,
            # Hybrid slides
            "agenda": self.add_agenda_slide,
            "timeline": self.add_timeline_slide,
            "icon-grid": self.add_icon_grid_slide,
            "stat-row": self.add_stat_row_slide,
            "pros-cons": self.add_pros_cons_slide,
            "comparison": self.add_comparison_slide,
            "checklist": self.add_checklist_slide,
            "logos": self.add_logos_slide,
            # Imported slides
            "architecture": self._import_catalog_slide,
        }

        for slide_data in slides:
            slide_type = slide_data.get("type", "content")
            if slide_type not in VALID_SLIDE_TYPES:
                print(f"Warning: Unknown slide type '{slide_type}', using 'content'")
                slide_type = "content"
            method = slide_methods.get(slide_type, self.add_content_slide)
            method(slide_data)

    def save(self, output_path: str) -> str:
        """Save presentation to file."""
        self.prs.save(output_path)
        return output_path


# =============================================================================
# CLI
# =============================================================================

def export_slides_as_images(pptx_path: str, output_dir: str) -> bool:
    """Export PPTX slides as PNG images using LibreOffice headless.

    Tries common LibreOffice installation paths on macOS, Linux, and
    falls back to the PATH. Returns True if successful.

    Args:
        pptx_path: Path to the .pptx file.
        output_dir: Directory to write PNG images to.
    """
    import subprocess
    import shutil

    soffice_candidates = [
        '/Applications/LibreOffice.app/Contents/MacOS/soffice',
        shutil.which('soffice'),
        shutil.which('libreoffice'),
        '/usr/bin/libreoffice',
        '/usr/bin/soffice',
    ]

    soffice = None
    for candidate in soffice_candidates:
        if candidate and Path(candidate).exists():
            soffice = candidate
            break

    if soffice is None:
        print("  LibreOffice not found — skipping image export.")
        print("  Install LibreOffice for visual verification, or review "
              "the .pptx manually.")
        return False

    Path(output_dir).mkdir(parents=True, exist_ok=True)

    try:
        result = subprocess.run(
            [
                soffice, '--headless', '--convert-to', 'png',
                '--outdir', output_dir, str(pptx_path),
            ],
            capture_output=True, text=True, timeout=120,
        )
        if result.returncode == 0:
            png_files = sorted(Path(output_dir).glob('*.png'))
            print(f"  Exported {len(png_files)} slide image(s) to {output_dir}/")
            for png in png_files:
                print(f"    {png.name}")
            return True
        else:
            print(f"  LibreOffice export failed: {result.stderr.strip()}")
            return False
    except subprocess.TimeoutExpired:
        print("  LibreOffice export timed out after 120 seconds")
        return False
    except Exception as e:
        print(f"  LibreOffice export error: {e}")
        return False


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Generate Databricks-branded PowerPoint presentations"
    )
    parser.add_argument("--input", "-i", required=True, help="Path to JSON content file")
    parser.add_argument("--output", "-o", required=True, help="Output path for .pptx file")
    parser.add_argument(
        "--verify", action="store_true",
        help="Export slides as PNG images (via LibreOffice) for visual verification"
    )

    args = parser.parse_args()

    # Load content
    with open(args.input, encoding='utf-8') as f:
        content = json.load(f)

    # Generate presentation
    generator = DatabricksSlideGenerator()
    generator.generate(content)
    output_path = generator.save(args.output)

    print(f"✓ Generated: {output_path}")
    print(f"  Slides: {generator.slide_count}")

    # Visual verification via image export
    if args.verify:
        verify_dir = str(Path(output_path).parent / "verify")
        print(f"\n  Exporting slides for visual verification...")
        export_slides_as_images(output_path, verify_dir)


if __name__ == "__main__":
    main()
