#!/usr/bin/env python3
"""
Build Icon Catalog from Architecture Catalog PPTX

Extracts icon-sized images from system_architecture_catalog.pptx,
deduplicates by SHA-256 hash, and writes each unique icon as a PNG
file plus an icon_catalog.json manifest.

Usage:
    python build-icon-catalog.py
"""

import hashlib
import json
import os
import re
import sys
import unicodedata
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("Error: python-pptx is required. Install with: pip3 install python-pptx")
    sys.exit(1)

try:
    from lxml import etree
except ImportError:
    print("Error: lxml is required. Install with: pip3 install lxml")
    sys.exit(1)

SCRIPT_DIR = Path(__file__).parent
SKILL_DIR = SCRIPT_DIR.parent
CATALOG_PPTX = SKILL_DIR / "assets" / "databricks" / "system_architecture_catalog.pptx"
OUTPUT_DIR = SKILL_DIR / "assets" / "icons"
OUTPUT_JSON = OUTPUT_DIR / "icon_catalog.json"

# ---------------------------------------------------------------------------
# Icon metadata: description, category, keywords, and optional renames
# ---------------------------------------------------------------------------
# Categories: aws, azure, gcp, databricks, data-platform, compute, security,
#             storage, data-format, ai-ml, integration, user-access, generic
#
# "rename" (optional) fixes misleading auto-derived names.
# Keys are the *current* auto-derived icon names (before rename).
# ---------------------------------------------------------------------------
ICON_METADATA = {
    # --- Databricks core ---
    "3p-data-platform": {
        "description": "Third-party data platform connector icon",
        "category": "integration",
        "keywords": ["third party", "connector", "external", "data platform"],
    },
    "account-console": {
        "description": "Databricks account console icon",
        "category": "databricks",
        "keywords": ["account", "console", "admin", "management"],
    },
    "admin": {
        "description": "Administrator/admin user icon",
        "category": "user-access",
        "keywords": ["admin", "administrator", "management", "user"],
    },
    "ai-models-tools": {
        "description": "Databricks AI models and tools icon",
        "category": "ai-ml",
        "keywords": ["ai", "models", "tools", "machine learning", "mosaic"],
    },
    "amazon-s3": {
        "description": "Amazon S3 logo (small, from overview slide)",
        "category": "aws",
        "keywords": ["aws", "s3", "storage", "amazon", "object storage"],
    },
    "anthropic": {
        "description": "Anthropic AI company logo",
        "category": "ai-ml",
        "keywords": ["anthropic", "claude", "ai", "llm", "foundation model"],
    },
    "apache-spark": {
        "description": "Apache Spark logo with text",
        "category": "data-platform",
        "keywords": ["spark", "apache", "distributed", "processing", "engine"],
    },
    "apps": {
        "description": "Databricks Apps icon",
        "category": "databricks",
        "keywords": ["apps", "applications", "databricks apps", "web apps"],
    },
    "athena": {
        "description": "Amazon Athena query service logo",
        "category": "aws",
        "keywords": ["athena", "aws", "query", "serverless", "sql"],
    },
    "audit-log": {
        "description": "Audit log / compliance icon",
        "category": "security",
        "keywords": ["audit", "log", "compliance", "tracking", "monitoring"],
    },
    # --- AWS downloaded icons ---
    "aws-dynamodb": {
        "description": "Amazon DynamoDB NoSQL database logo",
        "category": "aws",
        "keywords": ["aws", "dynamodb", "nosql", "database", "key-value"],
    },
    "aws-emr": {
        "description": "Amazon EMR (Elastic MapReduce) logo",
        "category": "aws",
        "keywords": ["aws", "emr", "mapreduce", "hadoop", "spark", "cluster"],
    },
    "aws-eventbridge": {
        "description": "Amazon EventBridge serverless event bus logo",
        "category": "aws",
        "keywords": ["aws", "eventbridge", "event", "serverless", "bus"],
    },
    "aws-glue": {
        "description": "AWS Glue ETL service logo",
        "category": "aws",
        "keywords": ["aws", "glue", "etl", "catalog", "crawl", "transform"],
    },
    "aws-iot-core": {
        "description": "AWS IoT Core logo",
        "category": "aws",
        "keywords": ["aws", "iot", "internet of things", "devices", "mqtt"],
    },
    "aws-kinesis": {
        "description": "Amazon Kinesis real-time streaming logo",
        "category": "aws",
        "keywords": ["aws", "kinesis", "streaming", "real-time", "data stream"],
    },
    "aws-lambda": {
        "description": "AWS Lambda serverless compute logo",
        "category": "aws",
        "keywords": ["aws", "lambda", "serverless", "function", "compute"],
    },
    "aws-rds": {
        "description": "Amazon RDS relational database logo",
        "category": "aws",
        "keywords": ["aws", "rds", "relational", "database", "managed"],
    },
    "aws-redshift": {
        "description": "Amazon Redshift data warehouse logo",
        "category": "aws",
        "keywords": ["aws", "redshift", "data warehouse", "analytics", "sql"],
    },
    "aws-s3": {
        "description": "Amazon S3 object storage logo",
        "category": "aws",
        "keywords": ["aws", "s3", "storage", "object storage", "bucket"],
    },
    "aws-sagemaker": {
        "description": "Amazon SageMaker ML platform logo",
        "category": "aws",
        "keywords": ["aws", "sagemaker", "machine learning", "ml", "training"],
    },
    # --- Azure downloaded icons ---
    "azure-active-directory": {
        "description": "Azure Active Directory (Entra ID) logo",
        "category": "azure",
        "keywords": ["azure", "active directory", "entra", "identity", "auth"],
    },
    "azure-cosmos-db": {
        "description": "Azure Cosmos DB multi-model database logo",
        "category": "azure",
        "keywords": ["azure", "cosmos", "nosql", "database", "global"],
    },
    "azure-data-factory": {
        "description": "Azure Data Factory ETL/ELT service logo",
        "category": "azure",
        "keywords": ["azure", "data factory", "etl", "pipeline", "orchestration"],
    },
    "azure-data-lake-storage": {
        "description": "Azure Data Lake Storage Gen2 logo",
        "category": "azure",
        "keywords": ["azure", "adls", "data lake", "storage", "gen2"],
    },
    "azure-event-hubs": {
        "description": "Azure Event Hubs streaming platform logo",
        "category": "azure",
        "keywords": ["azure", "event hubs", "streaming", "kafka", "ingestion"],
    },
    "azure-functions": {
        "description": "Azure Functions serverless compute logo",
        "category": "azure",
        "keywords": ["azure", "functions", "serverless", "compute", "faas"],
    },
    "azure-iot-hub": {
        "description": "Azure IoT Hub device management logo",
        "category": "azure",
        "keywords": ["azure", "iot", "hub", "devices", "telemetry"],
    },
    "azure-sql-database": {
        "description": "Azure SQL Database managed service logo",
        "category": "azure",
        "keywords": ["azure", "sql", "database", "managed", "relational"],
    },
    "azure-synapse": {
        "description": "Azure Synapse Analytics logo",
        "category": "azure",
        "keywords": ["azure", "synapse", "analytics", "data warehouse", "sql"],
    },
    # --- GCP downloaded icons ---
    "gcp-bigquery": {
        "description": "Google BigQuery serverless data warehouse logo",
        "category": "gcp",
        "keywords": ["gcp", "bigquery", "data warehouse", "analytics", "sql"],
    },
    "gcp-cloud-functions": {
        "description": "Google Cloud Functions serverless compute logo",
        "category": "gcp",
        "keywords": ["gcp", "cloud functions", "serverless", "faas", "compute"],
    },
    "gcp-cloud-sql": {
        "description": "Google Cloud SQL managed database logo",
        "category": "gcp",
        "keywords": ["gcp", "cloud sql", "database", "managed", "relational"],
    },
    "gcp-cloud-storage": {
        "description": "Google Cloud Storage object storage logo",
        "category": "gcp",
        "keywords": ["gcp", "cloud storage", "gcs", "object storage", "bucket"],
    },
    "gcp-dataflow": {
        "description": "Google Cloud Dataflow stream/batch processing logo",
        "category": "gcp",
        "keywords": ["gcp", "dataflow", "streaming", "batch", "beam", "processing"],
    },
    "gcp-dataproc": {
        "description": "Google Cloud Dataproc managed Spark/Hadoop logo",
        "category": "gcp",
        "keywords": ["gcp", "dataproc", "spark", "hadoop", "cluster", "managed"],
    },
    "gcp-pub-sub": {
        "description": "Google Cloud Pub/Sub messaging service logo",
        "category": "gcp",
        "keywords": ["gcp", "pub/sub", "messaging", "event", "streaming"],
    },
    "gcp-vertex-ai": {
        "description": "Google Vertex AI ML platform logo",
        "category": "gcp",
        "keywords": ["gcp", "vertex", "ai", "machine learning", "ml", "training"],
    },
    # --- Databricks platform components ---
    "batch-and-streaming": {
        "description": "Batch and streaming data processing icon",
        "category": "databricks",
        "keywords": ["batch", "streaming", "etl", "data processing", "pipeline"],
    },
    "bi": {
        "description": "Business intelligence / BI tools icon",
        "category": "databricks",
        "keywords": ["bi", "business intelligence", "dashboards", "reporting", "sql"],
    },
    "classic-compute": {
        "description": "Databricks classic compute cluster icon",
        "category": "compute",
        "keywords": ["classic", "compute", "cluster", "all-purpose", "jobs"],
    },
    "cloud-storage-1": {
        "rename": "delta-lake-logo",
        "description": "Delta Lake logo with text label",
        "category": "data-format",
        "keywords": ["delta lake", "delta", "lakehouse", "table format", "logo"],
    },
    "cloud-storage-2": {
        "rename": "cloud-database-large",
        "description": "Cloud database icon (large)",
        "category": "storage",
        "keywords": ["cloud", "database", "storage", "cloud database"],
    },
    "cloud-storage-3": {
        "rename": "cloud-database",
        "description": "Cloud database icon (small, used in reference architectures)",
        "category": "storage",
        "keywords": ["cloud", "database", "storage", "cloud database"],
    },
    "cluster": {
        "description": "Databricks cluster icon",
        "category": "compute",
        "keywords": ["cluster", "compute", "nodes", "spark", "workers"],
    },
    "cluster-disks": {
        "description": "Cluster with attached disks icon",
        "category": "compute",
        "keywords": ["cluster", "disks", "storage", "local", "ssd"],
    },
    "cluster-or-sql-warehouse": {
        "description": "Cluster or SQL warehouse icon (generic compute)",
        "category": "compute",
        "keywords": ["cluster", "sql warehouse", "compute", "endpoint"],
    },
    "clusters-1": {
        "rename": "cloud-cluster-green",
        "description": "Cloud cluster icon (green/serverless style)",
        "category": "compute",
        "keywords": ["cluster", "serverless", "cloud", "green", "compute"],
    },
    "clusters-2": {
        "rename": "cloud-cluster-orange",
        "description": "Cloud cluster icon (orange/classic style)",
        "category": "compute",
        "keywords": ["cluster", "classic", "cloud", "orange", "compute"],
    },
    "compute": {
        "description": "Databricks compute layer icon",
        "category": "compute",
        "keywords": ["compute", "processing", "layer", "databricks"],
    },
    "control-plane": {
        "description": "Databricks control plane icon",
        "category": "databricks",
        "keywords": ["control plane", "management", "orchestration", "api"],
    },
    "curated": {
        "description": "Curated/gold data layer icon (medallion architecture)",
        "category": "databricks",
        "keywords": ["curated", "gold", "medallion", "data quality", "refined"],
    },
    "data-shares": {
        "description": "Data sharing / Delta Sharing icon",
        "category": "databricks",
        "keywords": ["data sharing", "delta sharing", "share", "external"],
    },
    "database": {
        "description": "Generic database cylinder icon (small)",
        "category": "storage",
        "keywords": ["database", "cylinder", "storage", "relational"],
    },
    "delta-lake": {
        "description": "Delta Lake logo mark (triangular icon)",
        "category": "data-format",
        "keywords": ["delta lake", "delta", "lakehouse", "table format", "open source"],
    },
    "delta-sharing": {
        "description": "Delta Sharing protocol logo",
        "category": "databricks",
        "keywords": ["delta sharing", "sharing", "protocol", "open", "cross-platform"],
    },
    "devops-github-1": {
        "rename": "github-icon-1",
        "description": "GitHub-style icon (variant 1)",
        "category": "integration",
        "keywords": ["github", "git", "devops", "version control", "ci-cd"],
    },
    "devops-github-2": {
        "rename": "github-icon-2",
        "description": "GitHub-style icon (variant 2)",
        "category": "integration",
        "keywords": ["github", "git", "devops", "version control", "ci-cd"],
    },
    "dremio": {
        "description": "Dremio data lakehouse platform logo",
        "category": "data-platform",
        "keywords": ["dremio", "lakehouse", "query engine", "data lake"],
    },
    "driver": {
        "description": "Spark driver node icon",
        "category": "compute",
        "keywords": ["driver", "spark", "node", "master", "coordinator"],
    },
    "dwh-1": {
        "rename": "cloud-network",
        "description": "Cloud network topology icon",
        "category": "generic",
        "keywords": ["cloud", "network", "topology", "infrastructure", "nodes"],
    },
    "dwh-2": {
        "rename": "database-cylinder",
        "description": "Database cylinder icon",
        "category": "storage",
        "keywords": ["database", "cylinder", "dwh", "data warehouse", "storage"],
    },
    "emr": {
        "description": "Amazon EMR logo (from architecture slide)",
        "category": "aws",
        "keywords": ["emr", "aws", "elastic mapreduce", "hadoop", "spark"],
    },
    "enterprise-catalog": {
        "description": "Enterprise catalog / data catalog icon",
        "category": "databricks",
        "keywords": ["catalog", "enterprise", "metadata", "governance", "unity"],
    },
    "external-orchestrator": {
        "description": "External orchestrator / workflow tool icon",
        "category": "integration",
        "keywords": ["orchestrator", "workflow", "airflow", "external", "scheduler"],
    },
    "fabric": {
        "description": "Microsoft Fabric analytics platform logo",
        "category": "data-platform",
        "keywords": ["fabric", "microsoft", "analytics", "data platform"],
    },
    "feature-enhanced": {
        "description": "Feature engineering / enhanced features icon",
        "category": "ai-ml",
        "keywords": ["feature", "engineering", "enhanced", "ml", "transform"],
    },
    "feature-reduction": {
        "description": "Feature reduction / dimensionality reduction icon",
        "category": "ai-ml",
        "keywords": ["feature", "reduction", "dimensionality", "ml", "pca"],
    },
    "files-logs": {
        "description": "Files and logs data source icon",
        "category": "storage",
        "keywords": ["files", "logs", "data source", "unstructured", "raw"],
    },
    "flink": {
        "description": "Apache Flink stream processing logo",
        "category": "data-platform",
        "keywords": ["flink", "apache", "streaming", "processing", "real-time"],
    },
    "glue": {
        "description": "AWS Glue ETL logo (from architecture slide)",
        "category": "aws",
        "keywords": ["glue", "aws", "etl", "catalog", "transform"],
    },
    "google-cloud-storage": {
        "description": "Google Cloud Storage logo (small, from overview slide)",
        "category": "gcp",
        "keywords": ["gcp", "gcs", "cloud storage", "google", "object storage"],
    },
    "hms": {
        "description": "Hive Metastore (HMS) icon",
        "category": "data-platform",
        "keywords": ["hive", "metastore", "hms", "catalog", "schema"],
    },
    "huggingface": {
        "description": "Hugging Face ML community logo",
        "category": "ai-ml",
        "keywords": ["huggingface", "hugging face", "ml", "models", "transformers"],
    },
    "iceberg-1": {
        "rename": "iceberg-logo",
        "description": "Apache Iceberg logo with text",
        "category": "data-format",
        "keywords": ["iceberg", "apache", "table format", "open source", "logo"],
    },
    "iceberg-2": {
        "rename": "iceberg-mark",
        "description": "Apache Iceberg logo mark (small, no text)",
        "category": "data-format",
        "keywords": ["iceberg", "apache", "table format", "open source", "mark"],
    },
    "iceberg-rest": {
        "description": "Iceberg REST catalog icon",
        "category": "data-format",
        "keywords": ["iceberg", "rest", "catalog", "api", "open"],
    },
    "id-provider": {
        "description": "Identity provider / SSO icon",
        "category": "security",
        "keywords": ["identity", "provider", "sso", "auth", "oauth", "saml"],
    },
    "ingest-tool": {
        "description": "Data ingestion tool icon",
        "category": "databricks",
        "keywords": ["ingest", "ingestion", "tool", "lakeflow", "autoloader"],
    },
    "ingestion": {
        "description": "Data ingestion pipeline icon (medallion architecture)",
        "category": "databricks",
        "keywords": ["ingestion", "pipeline", "etl", "raw", "bronze"],
    },
    "jdbc-database": {
        "description": "JDBC database connector icon",
        "category": "integration",
        "keywords": ["jdbc", "database", "connector", "driver", "sql"],
    },
    "kafka": {
        "description": "Apache Kafka event streaming platform logo",
        "category": "data-platform",
        "keywords": ["kafka", "streaming", "event", "message queue", "pub-sub"],
    },
    "key-management": {
        "description": "Key management / encryption keys icon",
        "category": "security",
        "keywords": ["key management", "encryption", "kms", "cmk", "security"],
    },
    "lakeflowjobs": {
        "description": "Databricks Lakeflow Jobs logo/text",
        "category": "databricks",
        "keywords": ["lakeflow", "jobs", "orchestration", "workflow", "pipeline"],
    },
    "lakehouse": {
        "description": "Databricks Lakehouse icon",
        "category": "databricks",
        "keywords": ["lakehouse", "databricks", "platform", "unified", "architecture"],
    },
    "langchain": {
        "description": "LangChain LLM framework logo",
        "category": "ai-ml",
        "keywords": ["langchain", "llm", "framework", "chain", "agent", "rag"],
    },
    "managed-tables": {
        "description": "Managed tables icon (small)",
        "category": "databricks",
        "keywords": ["managed", "tables", "unity catalog", "governance"],
    },
    "marketplace": {
        "description": "Databricks Marketplace icon (small)",
        "category": "databricks",
        "keywords": ["marketplace", "data", "sharing", "listings", "exchange"],
    },
    "marketplaces": {
        "description": "Marketplaces / data exchange icon (reference arch)",
        "category": "databricks",
        "keywords": ["marketplace", "data exchange", "sharing", "listings"],
    },
    "media": {
        "description": "Media / unstructured content data source icon",
        "category": "generic",
        "keywords": ["media", "unstructured", "images", "video", "content"],
    },
    "microsoft-adls": {
        "description": "Microsoft Azure Data Lake Storage logo (small, overview)",
        "category": "azure",
        "keywords": ["azure", "adls", "data lake", "microsoft", "storage"],
    },
    "mlflow": {
        "description": "MLflow ML lifecycle management logo",
        "category": "ai-ml",
        "keywords": ["mlflow", "ml", "experiment", "tracking", "model registry"],
    },
    "notebooks": {
        "description": "Databricks notebooks icon",
        "category": "databricks",
        "keywords": ["notebook", "notebooks", "code", "interactive", "ide"],
    },
    "online-tables": {
        "description": "Databricks online tables icon (for model serving)",
        "category": "databricks",
        "keywords": ["online tables", "serving", "feature store", "low latency"],
    },
    "parquet": {
        "description": "Apache Parquet columnar format logo",
        "category": "data-format",
        "keywords": ["parquet", "columnar", "format", "apache", "file format"],
    },
    "photon-1": {
        "description": "Databricks Photon engine icon",
        "category": "databricks",
        "keywords": ["photon", "engine", "vectorized", "c++", "performance"],
    },
    "photon-2": {
        "rename": "apache-spark-logo",
        "description": "Apache Spark logo (small, horizontal)",
        "category": "data-platform",
        "keywords": ["spark", "apache", "logo", "distributed", "processing"],
    },
    "photon-3": {
        "rename": "apache-spark-logo-large",
        "description": "Apache Spark logo (large, horizontal)",
        "category": "data-platform",
        "keywords": ["spark", "apache", "logo", "distributed", "processing"],
    },
    "pyiceberg": {
        "description": "PyIceberg Python library icon",
        "category": "data-format",
        "keywords": ["pyiceberg", "python", "iceberg", "library", "api"],
    },
    "rdbms": {
        "description": "RDBMS / relational database source icon",
        "category": "storage",
        "keywords": ["rdbms", "relational", "database", "sql", "source"],
    },
    "redshift": {
        "description": "Amazon Redshift logo (from architecture slide)",
        "category": "aws",
        "keywords": ["redshift", "aws", "data warehouse", "analytics"],
    },
    "root-bucket": {
        "description": "Cloud root storage bucket icon",
        "category": "storage",
        "keywords": ["root", "bucket", "storage", "dbfs", "cloud"],
    },
    "salesforce": {
        "description": "Salesforce CRM platform logo",
        "category": "data-platform",
        "keywords": ["salesforce", "crm", "saas", "cloud", "data source"],
    },
    "secrets-1": {
        "rename": "key-vault",
        "description": "Key vault / secrets manager icon (large)",
        "category": "security",
        "keywords": ["key vault", "secrets", "credentials", "security", "vault"],
    },
    "secrets-2": {
        "rename": "key-small",
        "description": "Key icon (small)",
        "category": "security",
        "keywords": ["key", "secret", "credential", "token", "auth"],
    },
    "sensors-and-iot": {
        "description": "Sensors and IoT data source icon",
        "category": "generic",
        "keywords": ["sensors", "iot", "internet of things", "devices", "telemetry"],
    },
    "snowflake": {
        "description": "Snowflake data cloud platform logo",
        "category": "data-platform",
        "keywords": ["snowflake", "data warehouse", "cloud", "analytics"],
    },
    "spark-stream": {
        "description": "Spark Structured Streaming icon",
        "category": "databricks",
        "keywords": ["spark", "streaming", "structured streaming", "real-time"],
    },
    "sql-warehouses": {
        "description": "Databricks SQL Warehouses icon",
        "category": "compute",
        "keywords": ["sql warehouse", "endpoint", "serverless", "sql", "bi"],
    },
    "starburst": {
        "description": "Starburst (Trino enterprise) query engine logo",
        "category": "data-platform",
        "keywords": ["starburst", "trino", "query engine", "federation"],
    },
    "storage": {
        "description": "Databricks storage layer icon",
        "category": "storage",
        "keywords": ["storage", "layer", "cloud", "data", "databricks"],
    },
    "time-series-resampled-interpolated": {
        "description": "Time series resampling / interpolation icon",
        "category": "ai-ml",
        "keywords": ["time series", "resampling", "interpolation", "signal"],
    },
    "trino": {
        "description": "Trino distributed SQL query engine logo",
        "category": "data-platform",
        "keywords": ["trino", "distributed", "sql", "query engine", "presto"],
    },
    "unallocated-pool": {
        "description": "Unallocated resource pool icon",
        "category": "compute",
        "keywords": ["pool", "unallocated", "resources", "serverless", "shared"],
    },
    "unity-catalog": {
        "description": "Databricks Unity Catalog governance icon",
        "category": "databricks",
        "keywords": ["unity catalog", "governance", "metadata", "access control"],
    },
    "user": {
        "description": "Single user / person icon",
        "category": "user-access",
        "keywords": ["user", "person", "identity", "account"],
    },
    "username": {
        "description": "Username / credential badge icon",
        "category": "user-access",
        "keywords": ["username", "credential", "login", "identity", "badge"],
    },
    "users-1": {
        "rename": "users",
        "description": "Users / people group icon (large)",
        "category": "user-access",
        "keywords": ["users", "people", "group", "team", "access"],
    },
    "users-2": {
        "rename": "users-small",
        "description": "Users / people group icon (small)",
        "category": "user-access",
        "keywords": ["users", "people", "group", "team", "access"],
    },
    "workspace": {
        "description": "Databricks workspace icon",
        "category": "databricks",
        "keywords": ["workspace", "environment", "project", "databricks"],
    },
    "write-audit-log": {
        "description": "Write audit log / logging pipeline icon",
        "category": "security",
        "keywords": ["audit log", "write", "logging", "compliance", "trail"],
    },
}

# Icon size/file thresholds
MAX_DIMENSION_INCHES = 2.0   # Icons are < 2" on either axis
MAX_FILE_SIZE_BYTES = 200_000  # < 200KB

# EMU conversion
EMU_PER_INCH = 914400

# XML namespaces
NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'


def _emu_to_inches(emu):
    """Convert EMU to inches, rounded to 2 decimal places."""
    if emu is None:
        return 0.0
    return round(int(emu) / EMU_PER_INCH, 2)


def _normalize_name(text):
    """Derive a clean filename-safe icon name from text.

    Lowercase, strip parens/special chars, collapse whitespace to hyphens,
    truncate to 40 chars.
    """
    if not text:
        return ""
    # Normalize unicode
    text = unicodedata.normalize('NFKC', text)
    # Lowercase
    text = text.lower()
    # Strip parenthetical content
    text = re.sub(r'\([^)]*\)', '', text)
    # Replace special chars with space
    text = re.sub(r'[^a-z0-9\s-]', ' ', text)
    # Collapse whitespace to hyphens
    text = re.sub(r'\s+', '-', text.strip())
    # Remove leading/trailing hyphens
    text = text.strip('-')
    # Truncate
    return text[:40]


def _get_pic_position(pic_elem):
    """Extract position and size from a p:pic element's spPr/xfrm."""
    for ns in [NS_P, NS_A]:
        spPr = pic_elem.find(f'{{{ns}}}spPr')
        if spPr is not None:
            xfrm = spPr.find(f'{{{NS_A}}}xfrm')
            if xfrm is not None:
                off = xfrm.find(f'{{{NS_A}}}off')
                ext = xfrm.find(f'{{{NS_A}}}ext')
                if off is not None and ext is not None:
                    return {
                        "x_emu": int(off.get('x', '0')),
                        "y_emu": int(off.get('y', '0')),
                        "cx_emu": int(ext.get('cx', '0')),
                        "cy_emu": int(ext.get('cy', '0')),
                        "width": _emu_to_inches(ext.get('cx', '0')),
                        "height": _emu_to_inches(ext.get('cy', '0')),
                    }
    return None


def _get_pic_name(pic_elem):
    """Extract the name attribute from a p:pic element."""
    nvPicPr = pic_elem.find(f'{{{NS_P}}}nvPicPr')
    if nvPicPr is not None:
        cNvPr = nvPicPr.find(f'{{{NS_P}}}cNvPr')
        if cNvPr is not None:
            return cNvPr.get('name', '')
    return ''


def _get_group_transform(grpSp_elem):
    """Extract group coordinate transform."""
    for ns in [NS_P, NS_A]:
        grpSpPr = grpSp_elem.find(f'{{{ns}}}grpSpPr')
        if grpSpPr is not None:
            xfrm = grpSpPr.find(f'{{{NS_A}}}xfrm')
            if xfrm is not None:
                off = xfrm.find(f'{{{NS_A}}}off')
                ext = xfrm.find(f'{{{NS_A}}}ext')
                chOff = xfrm.find(f'{{{NS_A}}}chOff')
                chExt = xfrm.find(f'{{{NS_A}}}chExt')
                if all(e is not None for e in [off, ext, chOff, chExt]):
                    return {
                        "off_x": int(off.get('x', '0')),
                        "off_y": int(off.get('y', '0')),
                        "ext_cx": int(ext.get('cx', '0')),
                        "ext_cy": int(ext.get('cy', '0')),
                        "chOff_x": int(chOff.get('x', '0')),
                        "chOff_y": int(chOff.get('y', '0')),
                        "chExt_cx": int(chExt.get('cx', '0')),
                        "chExt_cy": int(chExt.get('cy', '0')),
                    }
    return None


def _child_to_slide_coords(child_x, child_y, gxfrm):
    """Convert child-space EMU coords to slide-space EMU."""
    if gxfrm["chExt_cx"] == 0 or gxfrm["chExt_cy"] == 0:
        return child_x, child_y
    slide_x = gxfrm["off_x"] + int(
        (child_x - gxfrm["chOff_x"]) * gxfrm["ext_cx"] / gxfrm["chExt_cx"]
    )
    slide_y = gxfrm["off_y"] + int(
        (child_y - gxfrm["chOff_y"]) * gxfrm["ext_cy"] / gxfrm["chExt_cy"]
    )
    return slide_x, slide_y


def _extract_shape_text(sp_elem):
    """Extract concatenated text from a shape element."""
    txBody = sp_elem.find(f'{{{NS_P}}}txBody')
    if txBody is None:
        txBody = sp_elem.find(f'{{{NS_A}}}txBody')
    if txBody is None:
        return ""
    texts = []
    for p in txBody.findall(f'{{{NS_A}}}p'):
        runs = p.findall(f'{{{NS_A}}}r')
        para_text = "".join(
            r.findtext(f'{{{NS_A}}}t', default='') for r in runs
        )
        if para_text.strip():
            texts.append(para_text.strip())
    return " ".join(texts)


def _find_nearest_text(pic_cx, pic_cy, text_shapes):
    """Find the nearest text shape to a pic's center, return its text.

    Uses Manhattan distance for simplicity.
    """
    best_dist = float('inf')
    best_text = ""
    for ts in text_shapes:
        dx = abs(ts["cx"] - pic_cx)
        dy = abs(ts["cy"] - pic_cy)
        dist = dx + dy
        if dist < best_dist:
            best_dist = dist
            best_text = ts["text"]
    return best_text


def _collect_text_shapes(parent, group_xfrm=None):
    """Collect text-bearing shapes with their center coordinates."""
    shapes = []
    for sp in parent.findall(f'{{{NS_P}}}sp'):
        text = _extract_shape_text(sp)
        if not text or len(text) <= 1:
            continue
        pos = _get_pic_position(sp)
        if pos:
            cx = pos["x_emu"] + pos["cx_emu"] // 2
            cy = pos["y_emu"] + pos["cy_emu"] // 2
            if group_xfrm:
                cx, cy = _child_to_slide_coords(cx, cy, group_xfrm)
            shapes.append({"text": text, "cx": cx, "cy": cy})

    for grpSp in parent.findall(f'{{{NS_P}}}grpSp'):
        gxfrm = _get_group_transform(grpSp)
        shapes.extend(_collect_text_shapes(grpSp, group_xfrm=gxfrm))

    return shapes


def _collect_pics(parent, slide, slide_num, group_xfrm=None):
    """Collect p:pic elements with their image data, position, and metadata."""
    pics = []
    for pic in parent.findall(f'{{{NS_P}}}pic'):
        pos = _get_pic_position(pic)
        if not pos:
            continue

        # Apply group transform if inside a group
        if group_xfrm:
            slide_x, slide_y = _child_to_slide_coords(
                pos["x_emu"], pos["y_emu"], group_xfrm
            )
            sx = group_xfrm["ext_cx"] / group_xfrm["chExt_cx"] if group_xfrm["chExt_cx"] else 1.0
            sy = group_xfrm["ext_cy"] / group_xfrm["chExt_cy"] if group_xfrm["chExt_cy"] else 1.0
            width = _emu_to_inches(int(pos["cx_emu"] * sx))
            height = _emu_to_inches(int(pos["cy_emu"] * sy))
            cx = slide_x + int(pos["cx_emu"] * sx) // 2
            cy = slide_y + int(pos["cy_emu"] * sy) // 2
        else:
            width = pos["width"]
            height = pos["height"]
            cx = pos["x_emu"] + pos["cx_emu"] // 2
            cy = pos["y_emu"] + pos["cy_emu"] // 2

        # Filter by icon size
        if width > MAX_DIMENSION_INCHES or height > MAX_DIMENSION_INCHES:
            continue

        # Get the image blob via blipFill
        blipFill = pic.find(f'{{{NS_P}}}blipFill')
        if blipFill is None:
            continue
        blip = blipFill.find(f'{{{NS_A}}}blip')
        if blip is None:
            continue
        rId = blip.get(f'{{{NS_R}}}embed')
        if not rId:
            continue

        try:
            rel = slide.part.rels[rId]
            image_part = rel.target_part
            blob = image_part.blob
            content_type = image_part.content_type
        except (KeyError, AttributeError):
            continue

        # Filter by file size
        if len(blob) > MAX_FILE_SIZE_BYTES:
            continue

        pic_name = _get_pic_name(pic)

        pics.append({
            "blob": blob,
            "content_type": content_type,
            "width": width,
            "height": height,
            "cx": cx,
            "cy": cy,
            "slide_num": slide_num,
            "pic_name": pic_name,
        })

    # Recurse into groups
    for grpSp in parent.findall(f'{{{NS_P}}}grpSp'):
        gxfrm = _get_group_transform(grpSp)
        pics.extend(_collect_pics(grpSp, slide, slide_num, group_xfrm=gxfrm))

    return pics


def build_icon_catalog():
    if not CATALOG_PPTX.exists():
        print(f"Error: Catalog not found at {CATALOG_PPTX}")
        sys.exit(1)

    prs = Presentation(str(CATALOG_PPTX))
    total_slides = len(prs.slides)
    print(f"Scanning {total_slides} slides for icons...")

    # Phase 1: Collect all icon-sized images with their positions
    all_pics = []
    all_text_shapes = {}  # slide_num -> list of text shapes

    for idx, slide in enumerate(prs.slides):
        slide_num = idx + 1
        cSld = slide._element.find(f'{{{NS_P}}}cSld')
        if cSld is None:
            continue
        spTree = cSld.find(f'{{{NS_P}}}spTree')
        if spTree is None:
            continue

        # Collect text shapes for nearest-text labeling
        text_shapes = _collect_text_shapes(spTree)
        all_text_shapes[slide_num] = text_shapes

        # Collect pic elements
        pics = _collect_pics(spTree, slide, slide_num)
        all_pics.extend(pics)

    print(f"  Found {len(all_pics)} icon-sized images across all slides")

    # Phase 2: Deduplicate by SHA-256 hash
    unique_icons = {}  # hash -> icon info
    for pic in all_pics:
        h = hashlib.sha256(pic["blob"]).hexdigest()
        if h not in unique_icons:
            # Find nearest text for labeling
            text_shapes = all_text_shapes.get(pic["slide_num"], [])
            nearest = _find_nearest_text(pic["cx"], pic["cy"], text_shapes)
            name = _normalize_name(nearest) or _normalize_name(pic["pic_name"])

            unique_icons[h] = {
                "blob": pic["blob"],
                "content_type": pic["content_type"],
                "width": pic["width"],
                "height": pic["height"],
                "nearest_text": nearest,
                "name": name,
                "source_slides": [pic["slide_num"]],
            }
        else:
            # Add slide to source_slides if not already there
            if pic["slide_num"] not in unique_icons[h]["source_slides"]:
                unique_icons[h]["source_slides"].append(pic["slide_num"])

    print(f"  Unique icons (by hash): {len(unique_icons)}")

    # Phase 3: Resolve name collisions
    name_counts = {}
    for h, icon in unique_icons.items():
        name = icon["name"]
        if not name:
            name = "icon"
            icon["name"] = name
        if name in name_counts:
            name_counts[name].append(h)
        else:
            name_counts[name] = [h]

    for name, hashes in name_counts.items():
        if len(hashes) > 1:
            for i, h in enumerate(hashes, start=1):
                unique_icons[h]["name"] = f"{name}-{i}"

    # Phase 4: Write icon files and build manifest
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    ext_map = {
        'image/png': '.png',
        'image/jpeg': '.jpg',
        'image/gif': '.gif',
        'image/svg+xml': '.svg',
        'image/x-emf': '.emf',
        'image/x-wmf': '.wmf',
        'image/tiff': '.tiff',
        'image/bmp': '.bmp',
    }

    # Build a reverse-rename map: old_name -> new_name
    rename_map = {}
    for old_name, meta in ICON_METADATA.items():
        if "rename" in meta:
            rename_map[old_name] = meta["rename"]

    catalog_icons = {}
    renamed_files = []

    for h, icon in sorted(unique_icons.items(), key=lambda x: x[1]["name"]):
        original_name = icon["name"]
        final_name = rename_map.get(original_name, original_name)

        ext = ext_map.get(icon["content_type"], '.png')
        filename = f'{final_name}{ext}'
        filepath = OUTPUT_DIR / filename

        with open(filepath, 'wb') as f:
            f.write(icon["blob"])

        # If renamed, remove old file (if it exists and differs)
        if final_name != original_name:
            old_filename = f'{original_name}{ext}'
            old_filepath = OUTPUT_DIR / old_filename
            if old_filepath.exists() and old_filepath != filepath:
                os.remove(old_filepath)
                renamed_files.append(f"  {original_name} -> {final_name}")

        # Merge metadata from ICON_METADATA
        meta = ICON_METADATA.get(original_name, {})
        entry = {
            "file": filename,
            "description": meta.get("description", ""),
            "category": meta.get("category", "generic"),
            "keywords": meta.get("keywords", []),
            "width": icon["width"],
            "height": icon["height"],
            "content_type": icon["content_type"],
            "source_slides": sorted(icon["source_slides"]),
        }
        catalog_icons[final_name] = entry

    # Phase 5: Preserve downloaded icons (source_slides: [])
    # These are cloud provider icons added manually, not extracted from PPTX.
    if OUTPUT_JSON.exists():
        with open(OUTPUT_JSON, 'r', encoding='utf-8') as f:
            existing_catalog = json.load(f)
        for name, entry in existing_catalog.get("icons", {}).items():
            if name not in catalog_icons and entry.get("source_slides") == []:
                # Downloaded icon — preserve it if the file still exists
                icon_file = OUTPUT_DIR / entry["file"]
                if icon_file.exists():
                    # Merge metadata if available (use original name or current name)
                    meta = ICON_METADATA.get(name, {})
                    entry["description"] = meta.get("description", entry.get("description", ""))
                    entry["category"] = meta.get("category", entry.get("category", "generic"))
                    entry["keywords"] = meta.get("keywords", entry.get("keywords", []))
                    catalog_icons[name] = entry

    # Sort icons by name for stable output
    sorted_icons = dict(sorted(catalog_icons.items()))

    catalog = {
        "source": CATALOG_PPTX.name,
        "icon_count": len(sorted_icons),
        "icons": sorted_icons,
    }

    with open(OUTPUT_JSON, 'w', encoding='utf-8') as f:
        json.dump(catalog, f, indent=2, ensure_ascii=False)

    print(f"\nWrote {len(sorted_icons)} icons to {OUTPUT_DIR}/")
    print(f"Wrote catalog to {OUTPUT_JSON}")

    if renamed_files:
        print(f"\nRenamed icons:")
        for r in renamed_files:
            print(r)

    # Check for icons without metadata
    missing_meta = [n for n in sorted_icons if n not in ICON_METADATA
                    and rename_map.get(n) not in ICON_METADATA
                    and not any(v == n for v in rename_map.values())]
    if missing_meta:
        print(f"\nWarning: {len(missing_meta)} icons without metadata:")
        for name in missing_meta:
            print(f"  {name}")

    # Print summary
    print(f"\nIcon names:")
    for name in sorted(sorted_icons.keys()):
        info = sorted_icons[name]
        desc = f" - {info['description']}" if info.get('description') else ""
        print(f"  {name} [{info['category']}]{desc}")


if __name__ == "__main__":
    build_icon_catalog()
